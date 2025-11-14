#!/usr/bin/env python3
"""
Generate minimal, scientific-style PowerPoint presentation.
Clean, understated design with programming aesthetic.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import numpy as np


# Minimal color scheme (grayscale + one accent)
COLORS = {
    'bg': RGBColor(255, 255, 255),      # White
    'text': RGBColor(33, 33, 33),       # Almost black
    'accent': RGBColor(94, 129, 172),   # Muted blue
    'light_gray': RGBColor(240, 240, 240),  # Very light gray
    'med_gray': RGBColor(180, 180, 180),    # Medium gray
    'dark_gray': RGBColor(100, 100, 100),   # Dark gray
}


def setup_presentation():
    """Create presentation with custom styling."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=COLORS['text'], align=PP_ALIGN.LEFT, mono=False):
    """Add a styled text box."""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align

    if mono:
        p.font.name = 'Consolas'
    else:
        p.font.name = 'Helvetica'

    return text_box


def create_title_slide(prs):
    """Minimal title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(1),
                 "MACE-Gaussian Interface",
                 font_size=36, bold=True, color=COLORS['text'], align=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.8),
                 "ML-Accelerated Anharmonic IR Spectroscopy",
                 font_size=18, color=COLORS['dark_gray'], align=PP_ALIGN.CENTER)

    # Bottom info
    add_text_box(slide, Inches(1), Inches(6.5), Inches(8), Inches(0.6),
                 "Your Name  ·  Research Group  ·  November 2025",
                 font_size=12, color=COLORS['med_gray'], align=PP_ALIGN.CENTER)

    # Subtle line
    line = slide.shapes.add_connector(1, Inches(3), Inches(4.5), Inches(7), Inches(4.5))
    line.line.color.rgb = COLORS['light_gray']
    line.line.width = Pt(1)

    return slide


def create_content_slide(prs, title, content_lines=None):
    """Minimal content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title with subtle underline
    add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.6),
                 title, font_size=24, bold=True, color=COLORS['text'])

    # Subtle line under title
    line = slide.shapes.add_connector(1, Inches(0.7), Inches(1.15), Inches(9.3), Inches(1.15))
    line.line.color.rgb = COLORS['light_gray']
    line.line.width = Pt(0.5)

    # Content
    if content_lines:
        y_pos = 1.6
        for level, text in content_lines:
            indent = level * 0.3
            bullet = "•" if level == 0 else "−"

            # Bullet
            add_text_box(slide, Inches(1.0 + indent), Inches(y_pos), Inches(0.2), Inches(0.3),
                        bullet, font_size=14, color=COLORS['accent'])

            # Text
            add_text_box(slide, Inches(1.3 + indent), Inches(y_pos), Inches(7.5 - indent), Inches(0.3),
                        text, font_size=14, color=COLORS['text'])

            y_pos += 0.45

    # Page number (bottom right)
    add_text_box(slide, Inches(9), Inches(7), Inches(0.5), Inches(0.3),
                str(len(prs.slides)), font_size=10, color=COLORS['med_gray'],
                align=PP_ALIGN.RIGHT)

    return slide


def create_image_slide(prs, title, image_path, caption=""):
    """Slide with single image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.6),
                 title, font_size=24, bold=True, color=COLORS['text'])

    # Line
    line = slide.shapes.add_connector(1, Inches(0.7), Inches(1.15), Inches(9.3), Inches(1.15))
    line.line.color.rgb = COLORS['light_gray']
    line.line.width = Pt(0.5)

    # Image
    try:
        pic = slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.8), width=Inches(7))
    except:
        # Placeholder if image doesn't exist
        add_text_box(slide, Inches(3), Inches(3.5), Inches(4), Inches(1),
                    f"[Image: {image_path}]", font_size=12, color=COLORS['med_gray'],
                    align=PP_ALIGN.CENTER)

    # Caption
    if caption:
        add_text_box(slide, Inches(1.5), Inches(6.5), Inches(7), Inches(0.4),
                    caption, font_size=11, color=COLORS['dark_gray'], align=PP_ALIGN.CENTER)

    # Page number
    add_text_box(slide, Inches(9), Inches(7), Inches(0.5), Inches(0.3),
                str(len(prs.slides)), font_size=10, color=COLORS['med_gray'],
                align=PP_ALIGN.RIGHT)

    return slide


def create_code_slide(prs, title, code_image_path):
    """Slide with code snippet."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.6),
                 title, font_size=24, bold=True, color=COLORS['text'])

    # Line
    line = slide.shapes.add_connector(1, Inches(0.7), Inches(1.15), Inches(9.3), Inches(1.15))
    line.line.color.rgb = COLORS['light_gray']
    line.line.width = Pt(0.5)

    # Code image
    try:
        pic = slide.shapes.add_picture(code_image_path, Inches(1), Inches(1.8), width=Inches(8))
    except:
        # Placeholder
        add_text_box(slide, Inches(2), Inches(3), Inches(6), Inches(2),
                    f"// Code snippet\n{code_image_path}", font_size=12,
                    color=COLORS['dark_gray'], align=PP_ALIGN.LEFT, mono=True)

    # Page number
    add_text_box(slide, Inches(9), Inches(7), Inches(0.5), Inches(0.3),
                str(len(prs.slides)), font_size=10, color=COLORS['med_gray'],
                align=PP_ALIGN.RIGHT)

    return slide


def create_two_column_slide(prs, title, left_lines, right_lines):
    """Two-column layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    add_text_box(slide, Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.6),
                 title, font_size=24, bold=True, color=COLORS['text'])

    # Line
    line = slide.shapes.add_connector(1, Inches(0.7), Inches(1.15), Inches(9.3), Inches(1.15))
    line.line.color.rgb = COLORS['light_gray']
    line.line.width = Pt(0.5)

    # Left column
    y_pos = 1.6
    for level, text in left_lines:
        indent = level * 0.3
        bullet = "•" if level == 0 else "−"

        add_text_box(slide, Inches(0.8 + indent), Inches(y_pos), Inches(0.2), Inches(0.3),
                    bullet, font_size=14, color=COLORS['accent'])
        add_text_box(slide, Inches(1.1 + indent), Inches(y_pos), Inches(3.5 - indent), Inches(0.3),
                    text, font_size=14, color=COLORS['text'])
        y_pos += 0.45

    # Vertical divider
    divider = slide.shapes.add_connector(1, Inches(5.0), Inches(1.4), Inches(5.0), Inches(7.0))
    divider.line.color.rgb = COLORS['light_gray']
    divider.line.width = Pt(0.5)

    # Right column
    y_pos = 1.6
    for level, text in right_lines:
        indent = level * 0.3
        bullet = "•" if level == 0 else "−"

        add_text_box(slide, Inches(5.3 + indent), Inches(y_pos), Inches(0.2), Inches(0.3),
                    bullet, font_size=14, color=COLORS['accent'])
        add_text_box(slide, Inches(5.6 + indent), Inches(y_pos), Inches(3.5 - indent), Inches(0.3),
                    text, font_size=14, color=COLORS['text'])
        y_pos += 0.45

    # Page number
    add_text_box(slide, Inches(9), Inches(7), Inches(0.5), Inches(0.3),
                str(len(prs.slides)), font_size=10, color=COLORS['med_gray'],
                align=PP_ALIGN.RIGHT)

    return slide


def create_section_slide(prs, title):
    """Minimal section divider."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Large title
    add_text_box(slide, Inches(1), Inches(3), Inches(8), Inches(1.5),
                 title, font_size=32, bold=True, color=COLORS['accent'],
                 align=PP_ALIGN.CENTER)

    return slide


def main():
    """Generate minimal scientific PowerPoint."""
    print("Creating minimal scientific PowerPoint...")
    print("=" * 60)

    prs = setup_presentation()

    # Slide 1: Title
    print("Creating slides...")
    create_title_slide(prs)

    # Slide 2: Motivation
    create_content_slide(prs, "Motivation", [
        (0, "Anharmonic spectroscopy is computationally expensive"),
        (1, "DFT calculations: hours to days per molecule"),
        (1, "Critical for accurate IR prediction (overtones, combinations)"),
        (0, "ML potentials offer orders of magnitude speedup"),
        (1, "Can we maintain accuracy while reducing cost?"),
    ])

    # Slide 3: Approach
    create_content_slide(prs, "Hybrid ML-QM Approach", [
        (0, "Use ML for fast components:"),
        (1, "Geometry optimization (MACE-MP/OMOL)"),
        (1, "Force constants (Hessian)"),
        (1, "Dipole moment derivatives"),
        (0, "Use Gaussian for anharmonic corrections:"),
        (1, "VPT2 perturbation theory"),
        (1, "Overtones and combination bands"),
        (0, "Result: 10-100× faster with comparable accuracy"),
    ])

    # Slide 4: Architecture
    create_image_slide(prs, "System Architecture",
                      "assets/flowcharts/workflow_terminal_minimal.png",
                      "End-to-end pipeline from XYZ input to HTML report")

    # Slide 5: ZMQ Bridge (two-column)
    create_two_column_slide(prs, "Technical Innovation: ZMQ Bridge",
        left_lines=[
            (0, "Challenge:"),
            (1, "Gaussian (Fortran) needs Python ML"),
            (1, "No direct FFI available"),
            (0, "Solution:"),
            (1, "Inter-process communication"),
            (1, "ZMQ message passing"),
            (1, "Unix socket IPC"),
        ],
        right_lines=[
            (0, "Mechanism:"),
            (1, "Main script: ZMQ server"),
            (1, "Helper script: ZMQ client"),
            (1, "Gaussian calls helper via External="),
            (1, "Real-time bidirectional data injection"),
        ]
    )

    # Slide 6: ZMQ Architecture
    create_image_slide(prs, "ZMQ Communication Flow",
                      "assets/flowcharts/zmq_terminal_minimal.png",
                      "Inter-process communication architecture")

    # Slide 7: Code example
    create_code_slide(prs, "Implementation: ZMQ Bridge",
                     "assets/code_snippets/pygments_material.png")

    # Slide 8: Module swapping
    create_image_slide(prs, "Module Swapping Mechanism",
                      "assets/flowcharts/module_swap_terminal_minimal.png",
                      "Dynamic switching between standard and dipole-enabled MACE")

    # Slide 9: Results
    create_image_slide(prs, "Results: Spectrum Comparison",
                      "assets/plots/sample_spectrum.png",
                      "ML vs DFT anharmonic IR spectra (example: water)")

    # Slide 10: Validation
    create_content_slide(prs, "Validation Strategy", [
        (0, "Rigorous mode-by-mode comparison"),
        (1, "Eigenvector-based mode matching"),
        (1, "Handles frequency errors and mode crossing"),
        (0, "Statistical metrics:"),
        (1, "MAE, RMSE, R² for frequencies"),
        (1, "Intensity correlations"),
        (0, "Multiple ML methods tested"),
        (1, "MACE-MP, MACE-OMOL"),
        (1, "Espaloma, custom MACE-ML dipoles"),
    ])

    # Slide 11: Impact
    create_section_slide(prs, "Impact & Future Work")

    # Slide 12: Conclusion
    create_content_slide(prs, "Summary", [
        (0, "Achievements:"),
        (1, "10-100× speedup over DFT anharmonic"),
        (1, "Novel ZMQ bridge for ML-QM integration"),
        (1, "Automated analysis workflow"),
        (0, "Future directions:"),
        (1, "Larger molecules (drugs, polymers)"),
        (1, "Improved dipole ML models"),
        (1, "High-throughput spectral database"),
    ])

    # Save
    output_file = "../presentation_minimal.pptx"
    prs.save(output_file)

    print("\n" + "=" * 60)
    print(f"Minimal PowerPoint created: {output_file}")
    print(f"Total slides: {len(prs.slides)}")
    print("\nStyle: Clean, scientific, minimal")
    print("Fonts: Helvetica (body), Consolas (code)")
    print("Colors: Grayscale + muted blue accent")


if __name__ == "__main__":
    main()
