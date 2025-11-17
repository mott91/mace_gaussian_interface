#!/usr/bin/env python3
"""
Test detailed workflow flowcharts with inputs, outputs, and side branches.
Multiple styles to choose from.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# Terminal Dark color scheme
BG = RGBColor(13, 17, 23)          # #0D1117 GitHub dark
TEXT = RGBColor(201, 209, 217)     # #C9D1D9
ACCENT = RGBColor(88, 166, 255)    # #58A6FF bright blue
DIM = RGBColor(139, 148, 158)      # #8B949E
GREEN = RGBColor(87, 171, 90)      # #57AB5A
ORANGE = RGBColor(219, 109, 40)    # #DB6D28
FONT = 'Consolas'


def create_detailed_vertical_workflow(prs):
    """Detailed vertical workflow with inputs/outputs and branches."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat workflow_detailed_v1.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Flowchart
    flow_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(6.2))
    tf = flow_box.text_frame

    lines = [
        ("INPUT: molecule.xyz (initial geometry)", GREEN),
        ("    │", DIM),
        ("    ▼", DIM),
        ("┌───────────────────────────────────────────────────────────────┐", ACCENT),
        ("│  Phase 1: Geometry Optimization                              │", ACCENT),
        ("│  Calculator: MACE-OMOL (default) or MACE-OFF/MACE-MP         │", TEXT),
        ("│  Output: optimized.xyz, results.json                         │", TEXT),
        ("└───────────────────────────────────────────────────────────────┘", ACCENT),
        ("    │", DIM),
        ("    ├──────────────────┐", DIM),
        ("    │                  │ (optional, default=True)", DIM),
        ("    ▼                  ▼", DIM),
        ("┌────────────────┐   ┌──────────────────────────────────────────┐", ACCENT),
        ("│  Phase 2:      │   │  Phase 3: ML Frequency Calculations      │", ACCENT),
        ("│  DFT Baseline  │   │  Energy: MACE-MP, MACE-OMOL              │", GREEN),
        ("│  B3LYP/        │   │  Dipole: Espaloma, MACE-ML               │", TEXT),
        ("│  6-31G(d,p)    │   │  ┌────────────────────────────────────┐  │", ORANGE),
        ("│  Output:       │   │  │ ZMQ Bridge                         │  │", ORANGE),
        ("│  reference     │   │  │ • Main: ZMQ server + ML dipoles    │  │", ORANGE),
        ("│  frequencies   │   │  │ • Helper: connects processes       │  │", ORANGE),
        ("└────────────────┘   │  │ • Gaussian: External calculation   │  │", ORANGE),
        ("    │                │  └────────────────────────────────────┘  │", ORANGE),
        ("    │                │  Output: {energy}_{dipole}/results.json  │", TEXT),
        ("    │                └──────────────────────────────────────────┘", ACCENT),
        ("    │                  │", DIM),
        ("    └──────────────────┘", DIM),
        ("    │", DIM),
        ("    ▼", DIM),
        ("┌───────────────────────────────────────────────────────────────┐", ACCENT),
        ("│  Phase 4: Analysis & Comparison                               │", ACCENT),
        ("│  • Compare ML vs DFT baseline                                 │", TEXT),
        ("│  • Metrics: MAE, RMSE, R², slope/intercept                    │", TEXT),
        ("│  • Generate regression plots                                  │", TEXT),
        ("│  • KDE broadening (FWHM=8 cm⁻¹)                               │", TEXT),
        ("│  Output: {molecule}_spectral_analysis.html                    │", TEXT),
        ("└───────────────────────────────────────────────────────────────┘", ACCENT),
    ]

    for i, (line, color) in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(9)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(1)
        p.line_spacing = 1.0


def create_swimlane_workflow(prs):
    """Swim lane style showing different components."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat workflow_swimlane.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Flowchart
    flow_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(9.2), Inches(6.2))
    tf = flow_box.text_frame

    lines = [
        ("═══════════════════════════════════════════════════════════════════", ACCENT),
        ("LANE 1: GEOMETRY    molecule.xyz → [MACE-OMOL] → optimized.xyz", ACCENT),
        ("═══════════════════════════════════════════════════════════════════", ACCENT),
        ("                                      │", DIM),
        ("═══════════════════════════════════════════════════════════════════", GREEN),
        ("LANE 2: DFT         optimized.xyz → [Gaussian B3LYP] → ref_freq", GREEN),
        ("═══════════════════════════════════════════════════════════════════", GREEN),
        ("                                      │", DIM),
        ("═══════════════════════════════════════════════════════════════════", ORANGE),
        ("LANE 3: ML FREQ     optimized.xyz ──┬→ [MACE-MP + Espaloma]", ORANGE),
        ("                                    ├→ [MACE-OMOL + Espaloma]", ORANGE),
        ("                                    ├→ [MACE-MP + MACE-ML]", ORANGE),
        ("                                    └→ [MACE-OMOL + MACE-ML]", ORANGE),
        ("", TEXT),
        ("                    ┌─────────────────────────────────────┐", ORANGE),
        ("                    │  ZMQ Bridge (Real-time)             │", ORANGE),
        ("                    │  Main ↔ Helper ↔ Gaussian           │", ORANGE),
        ("                    └─────────────────────────────────────┘", ORANGE),
        ("", TEXT),
        ("                    Each → {energy}_{dipole}/results.json", TEXT),
        ("═══════════════════════════════════════════════════════════════════", ORANGE),
        ("                                      │", DIM),
        ("═══════════════════════════════════════════════════════════════════", ACCENT),
        ("LANE 4: ANALYSIS    All results → [Compare] → HTML report", ACCENT),
        ("                    • Statistical metrics (MAE, RMSE, R²)", TEXT),
        ("                    • Regression plots", TEXT),
        ("                    • KDE spectra", TEXT),
        ("═══════════════════════════════════════════════════════════════════", ACCENT),
    ]

    for i, (line, color) in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(9)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(1)
        p.line_spacing = 1.0


def create_tree_workflow(prs):
    """Tree-style workflow showing parallel branches."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat workflow_tree.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Flowchart
    flow_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(6.2))
    tf = flow_box.text_frame

    lines = [
        ("molecule.xyz", GREEN),
        ("    │", DIM),
        ("    ├─▶ [Phase 1: Optimize with MACE-OMOL]", ACCENT),
        ("    │       → optimized.xyz", TEXT),
        ("    │", DIM),
        ("    ├─────┬─▶ [Phase 2: DFT Baseline]", GREEN),
        ("    │     │       B3LYP/6-31G(d,p)", TEXT),
        ("    │     │       → reference_frequencies.json", TEXT),
        ("    │     │", DIM),
        ("    │     └─▶ [Phase 3: ML Frequencies] ──┬─▶ MACE-MP + Espaloma", ORANGE),
        ("    │             Energy × Dipole         ├─▶ MACE-MP + MACE-ML", ORANGE),
        ("    │             combinations            ├─▶ MACE-OMOL + Espaloma", ORANGE),
        ("    │                                     └─▶ MACE-OMOL + MACE-ML", ORANGE),
        ("    │             │", DIM),
        ("    │             ├─▶ [ZMQ Bridge]", ORANGE),
        ("    │             │       Main Process (Python)", TEXT),
        ("    │             │       ├─▶ ZMQ Server", TEXT),
        ("    │             │       └─▶ ML Dipole Calc", TEXT),
        ("    │             │       │", DIM),
        ("    │             │       Helper Script", TEXT),
        ("    │             │       ├─▶ Connect Main ↔ Gaussian", TEXT),
        ("    │             │       └─▶ File path exchange", TEXT),
        ("    │             │       │", DIM),
        ("    │             │       Gaussian 16", TEXT),
        ("    │             │       └─▶ Frequency calculation", TEXT),
        ("    │             │", DIM),
        ("    │             └─▶ → {energy}_{dipole}/results.json", TEXT),
        ("    │", DIM),
        ("    └─────────────────▶ [Phase 4: Analysis]", ACCENT),
        ("                            Compare all methods", TEXT),
        ("                            ├─▶ Statistics (MAE, RMSE, R²)", TEXT),
        ("                            ├─▶ Regression plots", TEXT),
        ("                            ├─▶ KDE spectra (FWHM=8)", TEXT),
        ("                            └─▶ HTML report", TEXT),
    ]

    for i, (line, color) in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(10)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(1)
        p.line_spacing = 1.0


def create_compact_detailed_workflow(prs):
    """Compact but detailed workflow with annotations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat workflow_compact.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Flowchart
    flow_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.4), Inches(6.2))
    tf = flow_box.text_frame

    lines = [
        ("┌─────────────────────────────────────────────────────┐", GREEN),
        ("│  INPUT: molecule.xyz                                │", GREEN),
        ("└─────────────────────────────────────────────────────┘", GREEN),
        ("    ↓", DIM),
        ("┌─────────────────────────────────────────────────────┐", ACCENT),
        ("│  [1] Optimize → MACE-OMOL → optimized.xyz          │", ACCENT),
        ("└─────────────────────────────────────────────────────┘", ACCENT),
        ("    ↓", DIM),
        ("    ├──────────────────────┬─────────────────────────┐", DIM),
        ("    ▼                      ▼                         ▼", DIM),
        ("┌─────────────┐  ┌──────────────────────────────────────┐", GREEN),
        ("│  [2] DFT    │  │  [3] ML Frequencies (4 combos)       │", GREEN),
        ("│  Baseline   │  │  ┌────────────────────────────────┐  │", ORANGE),
        ("│  B3LYP      │  │  │  Energy: MACE-MP, MACE-OMOL    │  │", ORANGE),
        ("│  6-31G**    │  │  │  Dipole: Espaloma, MACE-ML     │  │", ORANGE),
        ("│             │  │  │  ────────────────────────────  │  │", TEXT),
        ("│  Output:    │  │  │  ZMQ: Main↔Helper↔Gaussian    │  │", ORANGE),
        ("│  reference  │  │  └────────────────────────────────┘  │", ORANGE),
        ("└─────────────┘  │  Output: 4 × results.json            │", TEXT),
        ("    │            └──────────────────────────────────────┘", GREEN),
        ("    │                      │", DIM),
        ("    └──────────────────────┘", DIM),
        ("    ↓", DIM),
        ("┌─────────────────────────────────────────────────────┐", ACCENT),
        ("│  [4] Compare & Analyze                              │", ACCENT),
        ("│  • Metrics: MAE, RMSE, R²                           │", TEXT),
        ("│  • Plots: regression, KDE spectra (FWHM=8)          │", TEXT),
        ("│  • Output: {molecule}_spectral_analysis.html        │", TEXT),
        ("└─────────────────────────────────────────────────────┘", ACCENT),
    ]

    for i, (line, color) in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(10)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(1)
        p.line_spacing = 1.0


def create_expanded_vertical_workflow(prs):
    """Expanded vertical workflow - maximum detail, full width."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat workflow_expanded.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Flowchart - use full width
    flow_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.95), Inches(9.4), Inches(6.35))
    tf = flow_box.text_frame

    lines = [
        ("INPUT: molecule.xyz (Cartesian coordinates, atom types)", GREEN),
        ("    │", DIM),
        ("    ▼", DIM),
        ("┌─────────────────────────────────────────────────────────────────────────────────┐", ACCENT),
        ("│  PHASE 1: GEOMETRY OPTIMIZATION                                                 │", ACCENT),
        ("│  Calculator: MACE-OMOL (default) | Options: MACE-OFF, MACE-MP                  │", TEXT),
        ("│  Method: BFGS optimizer, force convergence 0.01 eV/Å                           │", TEXT),
        ("│  Output: optimized.xyz (relaxed structure), geometry_opt/results.json          │", TEXT),
        ("│  Metrics: initial_energy, final_energy, n_steps, convergence_status            │", TEXT),
        ("└─────────────────────────────────────────────────────────────────────────────────┘", ACCENT),
        ("    │", DIM),
        ("    ├─────────────────────────────────────┬──────────────────────────────────────┐", DIM),
        ("    │                                     │                                      │", DIM),
        ("    │ (optional, --skip-dft-baseline)     │ (default: True)                      │", DIM),
        ("    ▼                                     ▼                                      ▼", DIM),
        ("┌───────────────────────┐  ┌──────────────────────────────────────────────────────────────┐", GREEN),
        ("│  PHASE 2: DFT         │  │  PHASE 3: ML FREQUENCY CALCULATIONS                          │", GREEN),
        ("│  BASELINE             │  │  Configuration: --energy-calculators, --dipole-calculators   │", ORANGE),
        ("│  ─────────────────    │  │  ─────────────────────────────────────────────────────────   │", TEXT),
        ("│  Method: B3LYP        │  │  Energy Calculators:                                         │", TEXT),
        ("│  Basis: 6-31G(d,p)    │  │    • MACE-MP (large, 2M params, general chemistry)           │", TEXT),
        ("│  Freq: Harmonic       │  │    • MACE-OMOL (molecules, 1.5M params, organic focus)       │", TEXT),
        ("│  Program: Gaussian 16 │  │  Dipole Calculators (with automatic fallback):               │", TEXT),
        ("│                       │  │    1. MACE-ML: Custom dipole model (~/dipole_model/)         │", ORANGE),
        ("│  Output:              │  │    2. Espaloma: Charge-based (Σ qᵢrᵢ)                        │", ORANGE),
        ("│  b3lyp_6-31Gdp/       │  │    3. xTB: GFN2-xTB semi-empirical                           │", ORANGE),
        ("│  results.json         │  │  ┌────────────────────────────────────────────────────────┐  │", ORANGE),
        ("│  ├─ frequencies[]     │  │  │  ZMQ BRIDGE ARCHITECTURE (IPC)                         │  │", ORANGE),
        ("│  ├─ intensities[]     │  │  │  ─────────────────────────────────                     │  │", ORANGE),
        ("│  └─ runtime_s         │  │  │  Main Process (gm_main.py):                            │  │", ORANGE),
        ("│                       │  │  │    • Binds ZMQ server socket (.ipc_file)               │  │", ORANGE),
        ("└───────────────────────┘  │  │    • Loads ML dipole calculator                        │  │", ORANGE),
        ("    │                      │  │    • Waits for geometry requests                       │  │", ORANGE),
        ("    │                      │  │    • Computes dipole derivatives                       │  │", ORANGE),
        ("    │                      │  │    • Writes Gaussian-format output                     │  │", ORANGE),
        ("    │                      │  │  Helper Script (gm_helper.py):                         │  │", ORANGE),
        ("    │                      │  │    • Launched by Gaussian External directive           │  │", ORANGE),
        ("    │                      │  │    • Connects to ZMQ socket                            │  │", ORANGE),
        ("    │                      │  │    • Sends: 'infile|outfile'                           │  │", ORANGE),
        ("    │                      │  │    • Waits for: 'done' signal                          │  │", ORANGE),
        ("    │                      │  │    • Returns control to Gaussian                       │  │", ORANGE),
        ("    │                      │  │  Gaussian 16:                                          │  │", ORANGE),
        ("    │                      │  │    • Freq calculation (harmonic/anharmonic)            │  │", ORANGE),
        ("    │                      │  │    • Calls external for each geometry                  │  │", ORANGE),
        ("    │                      │  └────────────────────────────────────────────────────────┘  │", ORANGE),
        ("    │                      │  Combinations tested (4 total):                              │", TEXT),
        ("    │                      │    • mace_mp_espaloma/results.json                           │", TEXT),
        ("    │                      │    • mace_mp_mace_ml/results.json                            │", TEXT),
        ("    │                      │    • mace_omol_espaloma/results.json                         │", TEXT),
        ("    │                      │    • mace_omol_mace_ml/results.json                          │", TEXT),
        ("    │                      └──────────────────────────────────────────────────────────────┘", GREEN),
        ("    │                                     │", DIM),
        ("    └─────────────────────────────────────┘", DIM),
        ("    │", DIM),
        ("    ▼", DIM),
        ("┌─────────────────────────────────────────────────────────────────────────────────┐", ACCENT),
        ("│  PHASE 4: ANALYSIS & COMPARISON (comparison_workflow.py)                        │", ACCENT),
        ("│  ─────────────────────────────────────────────────────────────────────────────  │", TEXT),
        ("│  Load Results: Parse all results.json files from comparison_results/            │", TEXT),
        ("│  Find Baseline: Auto-detect b3lyp_6-31Gdp as reference                          │", TEXT),
        ("│  Statistical Metrics:                                                           │", TEXT),
        ("│    • MAE (Mean Absolute Error) in cm⁻¹                                          │", TEXT),
        ("│    • RMSE (Root Mean Square Error) in cm⁻¹                                      │", TEXT),
        ("│    • R² (Coefficient of Determination)                                          │", TEXT),
        ("│    • Linear regression: slope, intercept                                        │", TEXT),
        ("│  Visualizations:                                                                │", TEXT),
        ("│    • Regression plots: ML freq vs DFT freq (scatter + fit line)                 │", TEXT),
        ("│    • KDE spectra: Gaussian broadening (FWHM=8.0 cm⁻¹, grid 400-4000 cm⁻¹)      │", TEXT),
        ("│    • Combined comparison plots (all methods overlaid)                           │", TEXT),
        ("│  Output: {molecule}_spectral_analysis.html (interactive, publication-ready)     │", TEXT),
        ("└─────────────────────────────────────────────────────────────────────────────────┘", ACCENT),
    ]

    for i, (line, color) in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(8)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(0)
        p.line_spacing = 1.0


def create_dim_only_workflow(prs):
    """All DIM (gray) version - understated monochrome style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat workflow_monochrome.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = DIM
    p.font.bold = True

    # Flowchart - use full width, all DIM
    flow_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.95), Inches(9.4), Inches(6.35))
    tf = flow_box.text_frame

    lines = [
        "INPUT: molecule.xyz (Cartesian coordinates, atom types)",
        "    │",
        "    ▼",
        "┌─────────────────────────────────────────────────────────────────────────────────┐",
        "│  PHASE 1: GEOMETRY OPTIMIZATION                                                 │",
        "│  Calculator: MACE-OMOL (default) | Options: MACE-OFF, MACE-MP                  │",
        "│  Method: BFGS optimizer, force convergence 0.01 eV/Å                           │",
        "│  Output: optimized.xyz (relaxed structure), geometry_opt/results.json          │",
        "│  Metrics: initial_energy, final_energy, n_steps, convergence_status            │",
        "└─────────────────────────────────────────────────────────────────────────────────┘",
        "    │",
        "    ├─────────────────────────────────────┬──────────────────────────────────────┐",
        "    │                                     │                                      │",
        "    │ (optional, --skip-dft-baseline)     │ (default: True)                      │",
        "    ▼                                     ▼                                      ▼",
        "┌───────────────────────┐  ┌──────────────────────────────────────────────────────────────┐",
        "│  PHASE 2: DFT         │  │  PHASE 3: ML FREQUENCY CALCULATIONS                          │",
        "│  BASELINE             │  │  Configuration: --energy-calculators, --dipole-calculators   │",
        "│  ─────────────────    │  │  ─────────────────────────────────────────────────────────   │",
        "│  Method: B3LYP        │  │  Energy Calculators:                                         │",
        "│  Basis: 6-31G(d,p)    │  │    • MACE-MP (large, 2M params, general chemistry)           │",
        "│  Freq: Harmonic       │  │    • MACE-OMOL (molecules, 1.5M params, organic focus)       │",
        "│  Program: Gaussian 16 │  │  Dipole Calculators (with automatic fallback):               │",
        "│                       │  │    1. MACE-ML: Custom dipole model (~/dipole_model/)         │",
        "│  Output:              │  │    2. Espaloma: Charge-based (Σ qᵢrᵢ)                        │",
        "│  b3lyp_6-31Gdp/       │  │    3. xTB: GFN2-xTB semi-empirical                           │",
        "│  results.json         │  │  ┌────────────────────────────────────────────────────────┐  │",
        "│  ├─ frequencies[]     │  │  │  ZMQ BRIDGE ARCHITECTURE (IPC)                         │  │",
        "│  ├─ intensities[]     │  │  │  ─────────────────────────────────                     │  │",
        "│  └─ runtime_s         │  │  │  Main Process (gm_main.py):                            │  │",
        "│                       │  │  │    • Binds ZMQ server socket (.ipc_file)               │  │",
        "└───────────────────────┘  │  │    • Loads ML dipole calculator                        │  │",
        "    │                      │  │    • Waits for geometry requests                       │  │",
        "    │                      │  │    • Computes dipole derivatives                       │  │",
        "    │                      │  │    • Writes Gaussian-format output                     │  │",
        "    │                      │  │  Helper Script (gm_helper.py):                         │  │",
        "    │                      │  │    • Launched by Gaussian External directive           │  │",
        "    │                      │  │    • Connects to ZMQ socket                            │  │",
        "    │                      │  │    • Sends: 'infile|outfile'                           │  │",
        "    │                      │  │    • Waits for: 'done' signal                          │  │",
        "    │                      │  │    • Returns control to Gaussian                       │  │",
        "    │                      │  │  Gaussian 16:                                          │  │",
        "    │                      │  │    • Freq calculation (harmonic/anharmonic)            │  │",
        "    │                      │  │    • Calls external for each geometry                  │  │",
        "    │                      │  └────────────────────────────────────────────────────────┘  │",
        "    │                      │  Combinations tested (4 total):                              │",
        "    │                      │    • mace_mp_espaloma/results.json                           │",
        "    │                      │    • mace_mp_mace_ml/results.json                            │",
        "    │                      │    • mace_omol_espaloma/results.json                         │",
        "    │                      │    • mace_omol_mace_ml/results.json                          │",
        "    │                      └──────────────────────────────────────────────────────────────┘",
        "    │                                     │",
        "    └─────────────────────────────────────┘",
        "    │",
        "    ▼",
        "┌─────────────────────────────────────────────────────────────────────────────────┐",
        "│  PHASE 4: ANALYSIS & COMPARISON (comparison_workflow.py)                        │",
        "│  ─────────────────────────────────────────────────────────────────────────────  │",
        "│  Load Results: Parse all results.json files from comparison_results/            │",
        "│  Find Baseline: Auto-detect b3lyp_6-31Gdp as reference                          │",
        "│  Statistical Metrics:                                                           │",
        "│    • MAE (Mean Absolute Error) in cm⁻¹                                          │",
        "│    • RMSE (Root Mean Square Error) in cm⁻¹                                      │",
        "│    • R² (Coefficient of Determination)                                          │",
        "│    • Linear regression: slope, intercept                                        │",
        "│  Visualizations:                                                                │",
        "│    • Regression plots: ML freq vs DFT freq (scatter + fit line)                 │",
        "│    • KDE spectra: Gaussian broadening (FWHM=8.0 cm⁻¹, grid 400-4000 cm⁻¹)      │",
        "│    • Combined comparison plots (all methods overlaid)                           │",
        "│  Output: {molecule}_spectral_analysis.html (interactive, publication-ready)     │",
        "└─────────────────────────────────────────────────────────────────────────────────┘",
    ]

    for i, line in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(8)
        p.font.name = FONT
        p.font.color.rgb = DIM
        p.space_after = Pt(0)
        p.line_spacing = 1.0


def main():
    """Generate detailed workflow test slides."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Generating detailed workflow diagrams...")

    create_detailed_vertical_workflow(prs)
    print("  ✓ Detailed vertical with branches")

    create_swimlane_workflow(prs)
    print("  ✓ Swim lane style")

    create_tree_workflow(prs)
    print("  ✓ Tree-style with parallel branches")

    create_compact_detailed_workflow(prs)
    print("  ✓ Compact detailed workflow")

    create_expanded_vertical_workflow(prs)
    print("  ✓ Expanded vertical (maximum detail, full width)")

    create_dim_only_workflow(prs)
    print("  ✓ DIM-only monochrome version (full width)")

    # Save presentation
    output_path = "../test_detailed_workflow.pptx"
    prs.save(output_path)
    print(f"\n✓ Test presentation created: {output_path}")
    print(f"  → 6 slides with detailed workflow variations")
    print(f"  → Shows inputs, outputs, branches, and ZMQ details")
    print(f"  → Terminal Dark styling")


if __name__ == "__main__":
    main()
