#!/usr/bin/env python3
"""
Generate complete 10-slide presentation using Terminal Dark template styling.
Based on the original template_terminal_dark.pptx design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# Terminal Dark color scheme
BG = RGBColor(13, 17, 23)          # #0D1117 GitHub dark
TEXT = RGBColor(201, 209, 217)     # #C9D1D9
ACCENT = RGBColor(88, 166, 255)    # #58A6FF bright blue
DIM = RGBColor(139, 148, 158)      # #8B949E
FONT = 'Consolas'


def add_footer(slide):
    """Add footer to slide (consistent across all slides)."""
    info_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.6))
    tf = info_box.text_frame
    tf.text = "Manuel Ott · Hofer Lab · 2025-11-19"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.name = FONT
    p.font.color.rgb = DIM
    p.alignment = PP_ALIGN.CENTER


def create_title_slide(prs):
    """Slide 0: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Prompt-style title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "$ ./presentation.sh\n> MACE-Gaussian Interface"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.6))
    tf = subtitle_box.text_frame
    tf.text = "// ML-accelerated IR spectroscopy"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.name = FONT
    p.font.color.rgb = DIM

    # Footer
    add_footer(slide)


def create_content_slide(prs, command, lines):
    """Create a content slide with command header and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title with prompt (28pt as requested)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = command
    p = tf.paragraphs[0]
    p.font.size = Pt(28)  # Increased from 20pt to 28pt
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(4.8))
    tf = content_box.text_frame

    for i, (line, color) in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(14)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(6)

    # Footer
    add_footer(slide)


def main():
    """Generate complete presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 0: Title
    create_title_slide(prs)

    # Slide 1: Motivation
    create_content_slide(prs, "$ cat motivation.md", [
        ("# Problem", ACCENT),
        ("  → DFT frequency calculations: slow", TEXT),
        ("  → Minutes to hours per molecule", TEXT),
        ("  → Bottleneck for high-throughput screening", TEXT),
        ("", TEXT),
        ("# Solution", ACCENT),
        ("  → Hybrid ML-QM approach", TEXT),
        ("  → ML dipole derivatives via ZMQ bridge", TEXT),
        ("  → 10-100× faster than pure DFT", TEXT),
        ("", TEXT),
        ("# Impact", ACCENT),
        ("  → Enable rapid IR spectral predictions", TEXT),
        ("  → Maintain DFT-level accuracy", TEXT),
    ])

    # Slide 2: Architecture Overview
    create_content_slide(prs, "$ cat architecture.md", [
        ("# System Architecture", ACCENT),
        ("", TEXT),
        ("┌─────────────┐      ┌──────────────┐", DIM),
        ("│ MACE        │──────│ Gaussian 16  │", DIM),
        ("│ (Python)    │ ZMQ  │ (Fortran)    │", DIM),
        ("└─────────────┘      └──────────────┘", DIM),
        ("", TEXT),
        ("# Workflow Phases", ACCENT),
        ("  1. Geometry optimization (MACE-OMOL)", TEXT),
        ("  2. DFT baseline (B3LYP/6-31G(d,p))", TEXT),
        ("  3. ML frequency calculations", TEXT),
        ("  4. Statistical comparison & analysis", TEXT),
    ])

    # Slide 3: ZMQ Bridge
    create_content_slide(prs, "$ cat zmq_bridge.md", [
        ("# Inter-Process Communication", ACCENT),
        ("", TEXT),
        ("  1. Gaussian calls external program", TEXT),
        ("  2. Helper script connects via ZMQ socket", TEXT),
        ("  3. Main process receives geometry", TEXT),
        ("  4. ML calculator computes dipoles", TEXT),
        ("  5. Results written to Gaussian format", TEXT),
        ("  6. Gaussian continues calculation", TEXT),
        ("", TEXT),
        ("# Key Innovation", ACCENT),
        ("  → Real-time injection of ML predictions", TEXT),
        ("  → No modification to Gaussian source", TEXT),
        ("  → File-based IPC for reliability", TEXT),
    ])

    # Slide 4: Dipole Calculators
    create_content_slide(prs, "$ ls -la calculators/", [
        ("# Modular Dipole System", ACCENT),
        ("", TEXT),
        ("drwxr-xr-x  mace_ml/          # Custom MACE", DIM),
        ("drwxr-xr-x  espaloma/         # ML charges", DIM),
        ("drwxr-xr-x  xtb/              # Semi-empirical", DIM),
        ("", TEXT),
        ("# Automatic Fallback Chain", ACCENT),
        ("  1. MACE-ML: Custom dipole model", TEXT),
        ("  2. Espaloma: Charge-based dipoles", TEXT),
        ("  3. xTB: GFN2-xTB semi-empirical", TEXT),
        ("", TEXT),
        ("# Implementation", ACCENT),
        ("  → Factory pattern with availability checks", TEXT),
        ("  → Seamless calculator switching", TEXT),
    ])

    # Slide 5: Module Swapping
    create_content_slide(prs, "$ cat module_swap.py", [
        ("# Dual MACE Architecture", ACCENT),
        ("", TEXT),
        ("mace_ML_pkg/          # Standard MACE", DIM),
        ("mace_dipole_pkg/      # Custom fork", DIM),
        ("", TEXT),
        ("# Implementation", ACCENT),
        ("  → Dynamic module replacement", TEXT),
        ("  → Cache invalidation required", TEXT),
        ("  → Cleanup after dipole calculations", TEXT),
        ("", TEXT),
        ("# Challenge", ACCENT),
        ("  → Prevent module contamination", TEXT),
        ("  → Thread-safe switching mechanism", TEXT),
    ])

    # Slide 6: Validation Method
    create_content_slide(prs, "$ cat validation.md", [
        ("# Comparison vs DFT Baseline", ACCENT),
        ("", TEXT),
        ("  → B3LYP/6-31G(d,p) reference", TEXT),
        ("  → Harmonic frequency analysis", TEXT),
        ("  → Test molecule: Water (H₂O)", TEXT),
        ("", TEXT),
        ("# Statistical Metrics", ACCENT),
        ("  → MAE (Mean Absolute Error)", TEXT),
        ("  → RMSE (Root Mean Square Error)", TEXT),
        ("  → R² (Coefficient of Determination)", TEXT),
        ("  → Linear regression analysis", TEXT),
        ("", TEXT),
        ("# Visualization", ACCENT),
        ("  → Regression plots (ML vs DFT)", TEXT),
        ("  → Gaussian KDE broadening (8 cm⁻¹)", TEXT),
    ])

    # Slide 7: Results - Water
    create_content_slide(prs, "$ python analyze_spectra.py water", [
        ("# Water (H₂O) Benchmark", ACCENT),
        ("", TEXT),
        ("  → 3 normal modes", TEXT),
        ("  → Symmetric stretch (~3650 cm⁻¹)", TEXT),
        ("  → Asymmetric stretch (~3750 cm⁻¹)", TEXT),
        ("  → Bending mode (~1590 cm⁻¹)", TEXT),
        ("", TEXT),
        ("# ML vs DFT Accuracy", ACCENT),
        ("  → MAE: < 20 cm⁻¹", TEXT),
        ("  → RMSE: < 25 cm⁻¹", TEXT),
        ("  → R²: > 0.99", TEXT),
        ("", TEXT),
        ("# Computational Speedup", ACCENT),
        ("  → 50-100× faster than pure DFT", TEXT),
    ])

    # Slide 8: Calculator Combinations
    create_content_slide(prs, "$ cat config.yaml", [
        ("# Energy Calculators", ACCENT),
        ("", TEXT),
        ("  → MACE-MP (large, general)", TEXT),
        ("  → MACE-OMOL (molecules, default)", TEXT),
        ("  → MACE-OFF (alternative foundation)", TEXT),
        ("", TEXT),
        ("# Dipole Calculators", ACCENT),
        ("  → Espaloma (charge-based)", TEXT),
        ("  → MACE-ML (custom dipole model)", TEXT),
        ("", TEXT),
        ("# Tested Combinations", ACCENT),
        ("  → mace_mp + espaloma", TEXT),
        ("  → mace_omol + espaloma", TEXT),
        ("  → mace_mp + mace_ml", TEXT),
        ("  → All compared vs DFT baseline", TEXT),
    ])

    # Slide 9: Workflow Pipeline
    create_content_slide(prs, "$ python cli.py run water.xyz", [
        ("# 4-Phase Pipeline", ACCENT),
        ("", TEXT),
        ("  Phase 1: Geometry optimization", TEXT),
        ("    → MACE-OMOL calculator", TEXT),
        ("", TEXT),
        ("  Phase 2: DFT baseline", TEXT),
        ("    → B3LYP/6-31G(d,p) frequencies", TEXT),
        ("", TEXT),
        ("  Phase 3: ML frequency calculations", TEXT),
        ("    → Multiple energy-dipole combinations", TEXT),
        ("", TEXT),
        ("  Phase 4: Analysis & comparison", TEXT),
        ("    → Statistics, plots, HTML report", TEXT),
    ])

    # Slide 10: Future Directions
    create_content_slide(prs, "$ git log --oneline --graph", [
        ("# Current Status", ACCENT),
        ("  → ZMQ bridge implemented & working", TEXT),
        ("  → Water benchmark validated", TEXT),
        ("  → Automated 4-phase pipeline", TEXT),
        ("  → Multiple calculator combinations", TEXT),
        ("", TEXT),
        ("# Next Steps", ACCENT),
        ("  → Expand to larger molecules", TEXT),
        ("  → Anharmonic frequency analysis", TEXT),
        ("  → Optimize MACE-ML dipole model", TEXT),
        ("  → ORCA support (beyond Gaussian)", TEXT),
        ("", TEXT),
        ("# Long-term Vision", ACCENT),
        ("  → Universal ML-QM spectroscopy platform", TEXT),
    ])

    # Save presentation
    output_path = "../mace_gaussian_presentation.pptx"
    prs.save(output_path)
    print(f"✓ Full presentation created: {output_path}")
    print(f"  → 1 title slide + 10 content slides")
    print(f"  → Headers: 28pt")
    print(f"  → Footer on all slides")
    print(f"  → Terminal Dark styling")


if __name__ == "__main__":
    main()
