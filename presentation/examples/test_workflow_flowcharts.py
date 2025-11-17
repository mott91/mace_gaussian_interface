#!/usr/bin/env python3
"""
Test different ASCII flowchart layouts for the 4-phase workflow.
Generate multiple slides with various flowchart styles.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# Terminal Dark color scheme
BG = RGBColor(13, 17, 23)          # #0D1117 GitHub dark
TEXT = RGBColor(201, 209, 217)     # #C9D1D9
ACCENT = RGBColor(88, 166, 255)    # #58A6FF bright blue
DIM = RGBColor(139, 148, 158)      # #8B949E
FONT = 'Consolas'


def create_vertical_workflow_slide(prs):
    """Vertical flowchart - top to bottom."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat workflow_vertical.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Flowchart box
    flow_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(7), Inches(5.8))
    tf = flow_box.text_frame

    flowchart_lines = [
        ("water.xyz", TEXT),
        ("    │", DIM),
        ("    ▼", DIM),
        ("┌──────────────────────────────────────────┐", ACCENT),
        ("│  Phase 1: Geometry Optimization         │", ACCENT),
        ("│  → MACE-OMOL calculator                  │", TEXT),
        ("│  → Output: optimized.xyz                 │", TEXT),
        ("└──────────────────────────────────────────┘", ACCENT),
        ("    │", DIM),
        ("    ▼", DIM),
        ("┌──────────────────────────────────────────┐", ACCENT),
        ("│  Phase 2: DFT Baseline                   │", ACCENT),
        ("│  → B3LYP/6-31G(d,p)                      │", TEXT),
        ("│  → Output: reference frequencies         │", TEXT),
        ("└──────────────────────────────────────────┘", ACCENT),
        ("    │", DIM),
        ("    ▼", DIM),
        ("┌──────────────────────────────────────────┐", ACCENT),
        ("│  Phase 3: ML Frequency Calculations      │", ACCENT),
        ("│  → Energy: MACE-MP/MACE-OMOL             │", TEXT),
        ("│  → Dipole: Espaloma/MACE-ML              │", TEXT),
        ("│  → ZMQ bridge → Gaussian                 │", TEXT),
        ("└──────────────────────────────────────────┘", ACCENT),
        ("    │", DIM),
        ("    ▼", DIM),
        ("┌──────────────────────────────────────────┐", ACCENT),
        ("│  Phase 4: Analysis & Comparison          │", ACCENT),
        ("│  → Statistics: MAE, RMSE, R²             │", TEXT),
        ("│  → Plots: regression, spectra            │", TEXT),
        ("│  → Output: HTML report                   │", TEXT),
        ("└──────────────────────────────────────────┘", ACCENT),
    ]

    for i, (line, color) in enumerate(flowchart_lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(11)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(2)
        p.line_spacing = 1.0


def create_horizontal_workflow_slide(prs):
    """Horizontal flowchart - left to right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat workflow_horizontal.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Flowchart box
    flow_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = flow_box.text_frame

    flowchart_lines = [
        ("┌─────────┐        ┌─────────┐        ┌─────────┐        ┌─────────┐", ACCENT),
        ("│ Phase 1 │───────▶│ Phase 2 │───────▶│ Phase 3 │───────▶│ Phase 4 │", ACCENT),
        ("└─────────┘        └─────────┘        └─────────┘        └─────────┘", ACCENT),
        ("", TEXT),
        ("", TEXT),
        ("Geometry           DFT                ML                 Analysis", ACCENT),
        ("Optimization       Baseline           Frequencies        & Reports", ACCENT),
        ("", TEXT),
        ("", TEXT),
        ("MACE-OMOL          B3LYP              MACE-MP +          Statistics", TEXT),
        ("calculator         6-31G(d,p)         Espaloma           Regression", TEXT),
        ("", TEXT),
        ("optimized.xyz      reference          ZMQ bridge         HTML", TEXT),
        ("                   frequencies        Gaussian           plots", TEXT),
    ]

    for i, (line, color) in enumerate(flowchart_lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(13)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(4)


def create_detailed_workflow_slide(prs):
    """Detailed flowchart with more information."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat workflow_detailed.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Flowchart box
    flow_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(8.4), Inches(5.9))
    tf = flow_box.text_frame

    flowchart_lines = [
        ("INPUT: water.xyz", ACCENT),
        ("  ↓", DIM),
        ("╔═══════════════════════════════════════════════════════════╗", ACCENT),
        ("║  PHASE 1: GEOMETRY OPTIMIZATION                           ║", ACCENT),
        ("╠═══════════════════════════════════════════════════════════╣", ACCENT),
        ("║  Calculator: MACE-OMOL (default)                          ║", TEXT),
        ("║  Output:     optimized.xyz, results.json                  ║", TEXT),
        ("╚═══════════════════════════════════════════════════════════╝", ACCENT),
        ("  ↓", DIM),
        ("╔═══════════════════════════════════════════════════════════╗", ACCENT),
        ("║  PHASE 2: DFT BASELINE (optional, default=True)           ║", ACCENT),
        ("╠═══════════════════════════════════════════════════════════╣", ACCENT),
        ("║  Method:     B3LYP/6-31G(d,p)                             ║", TEXT),
        ("║  Type:       Harmonic frequencies                         ║", TEXT),
        ("║  Output:     b3lyp_6-31Gdp/results.json                   ║", TEXT),
        ("╚═══════════════════════════════════════════════════════════╝", ACCENT),
        ("  ↓", DIM),
        ("╔═══════════════════════════════════════════════════════════╗", ACCENT),
        ("║  PHASE 3: ML FREQUENCY CALCULATIONS                       ║", ACCENT),
        ("╠═══════════════════════════════════════════════════════════╣", ACCENT),
        ("║  Energy:     MACE-MP, MACE-OMOL                           ║", TEXT),
        ("║  Dipole:     Espaloma, MACE-ML                            ║", TEXT),
        ("║  Bridge:     ZMQ IPC → Gaussian 16                        ║", TEXT),
        ("║  Output:     {energy}_{dipole}/results.json               ║", TEXT),
        ("╚═══════════════════════════════════════════════════════════╝", ACCENT),
        ("  ↓", DIM),
        ("╔═══════════════════════════════════════════════════════════╗", ACCENT),
        ("║  PHASE 4: ANALYSIS & COMPARISON                           ║", ACCENT),
        ("╠═══════════════════════════════════════════════════════════╣", ACCENT),
        ("║  Metrics:    MAE, RMSE, R²                                ║", TEXT),
        ("║  Plots:      Regression, KDE spectra                      ║", TEXT),
        ("║  Output:     {molecule}_spectral_analysis.html            ║", TEXT),
        ("╚═══════════════════════════════════════════════════════════╝", ACCENT),
    ]

    for i, (line, color) in enumerate(flowchart_lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(10)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(1)
        p.line_spacing = 1.0


def create_compact_workflow_slide(prs):
    """Compact flowchart - minimalist style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat workflow_compact.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Flowchart box
    flow_box = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(6), Inches(4.5))
    tf = flow_box.text_frame

    flowchart_lines = [
        ("water.xyz", TEXT),
        ("    │", DIM),
        ("    ├─▶ [1] Optimize (MACE-OMOL)", ACCENT),
        ("    │", DIM),
        ("    ├─▶ [2] DFT Baseline (B3LYP)", ACCENT),
        ("    │", DIM),
        ("    ├─▶ [3] ML Freq (Gaussian+ZMQ)", ACCENT),
        ("    │       Energy: MACE-MP", TEXT),
        ("    │       Dipole: Espaloma", TEXT),
        ("    │", DIM),
        ("    └─▶ [4] Analysis (MAE/RMSE/R²)", ACCENT),
        ("", TEXT),
        ("", TEXT),
        ("OUTPUT: HTML report", TEXT),
    ]

    for i, (line, color) in enumerate(flowchart_lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(14)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(6)


def main():
    """Generate test presentation with different flowchart styles."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Generating workflow flowchart test slides...")

    # Create different flowchart styles
    create_vertical_workflow_slide(prs)
    print("  ✓ Vertical flowchart")

    create_horizontal_workflow_slide(prs)
    print("  ✓ Horizontal flowchart")

    create_detailed_workflow_slide(prs)
    print("  ✓ Detailed flowchart")

    create_compact_workflow_slide(prs)
    print("  ✓ Compact flowchart")

    # Save presentation
    output_path = "../test_workflow_flowcharts.pptx"
    prs.save(output_path)
    print(f"\n✓ Test presentation created: {output_path}")
    print(f"  → 4 slides with different flowchart styles")
    print(f"  → Terminal Dark styling")
    print(f"\nReview the slides and pick your favorite style!")


if __name__ == "__main__":
    main()
