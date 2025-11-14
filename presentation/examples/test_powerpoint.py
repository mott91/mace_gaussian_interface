#!/usr/bin/env python3
"""
Test PowerPoint generation using python-pptx.

This script demonstrates:
1. Creating slides programmatically
2. Adding text with formatting
3. Embedding images
4. Creating tables
5. Custom styling and themes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np


def create_title_slide(prs, title, subtitle):
    """Create a title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Subtitle
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)

    return slide


def create_section_header(prs, title):
    """Create a section header slide."""
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def create_content_slide(prs, title, bullet_points=None):
    """Create a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Content
    if bullet_points:
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()

        for level, text in bullet_points:
            p = text_frame.add_paragraph()
            p.text = text
            p.level = level
            p.font.size = Pt(18)

    return slide


def create_image_slide(prs, title, image_path, caption=""):
    """Create a slide with an image."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add image
    left = Inches(1.5)
    top = Inches(1.5)
    width = Inches(7)
    pic = slide.shapes.add_picture(image_path, left, top, width=width)

    # Add caption if provided
    if caption:
        left = Inches(1.5)
        top = Inches(6.5)
        width = Inches(7)
        height = Inches(0.5)
        caption_box = slide.shapes.add_textbox(left, top, width, height)
        caption_frame = caption_box.text_frame
        caption_frame.text = caption
        p = caption_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide


def create_two_column_slide(prs, title, left_content, right_content):
    """Create a slide with two columns."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True

    # Left column
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(4.5)
    height = Inches(5)
    left_box = slide.shapes.add_textbox(left, top, width, height)
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for level, text in left_content:
        p = left_frame.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(16)

    # Right column
    left = Inches(5.2)
    right_box = slide.shapes.add_textbox(left, top, width, height)
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for level, text in right_content:
        p = right_frame.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(16)

    return slide


def create_code_slide(prs, title, code_image_path):
    """Create a slide with a code snippet image."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True

    # Add code image (terminal style)
    left = Inches(0.8)
    top = Inches(1.4)
    width = Inches(8.4)
    try:
        pic = slide.shapes.add_picture(code_image_path, left, top, width=width)
    except Exception as e:
        # If image doesn't exist, add placeholder
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = f"[Code snippet image placeholder]\n{code_image_path}"
        p = text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

    return slide


def create_flowchart_slide(prs, title, flowchart_image_path):
    """Create a slide with a flowchart."""
    return create_image_slide(prs, title, flowchart_image_path)


def add_styled_textbox(slide, left, top, width, height, text, font_size=18,
                       bold=False, color=RGBColor(0, 0, 0), bg_color=None):
    """Helper to add a styled text box."""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color

    if bg_color:
        fill = text_box.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    return text_box


def create_sample_plot():
    """Create a sample plot for demonstration."""
    fig, ax = plt.subplots(figsize=(8, 5), dpi=150)

    # Sample data
    x = np.linspace(0, 4000, 1000)
    y_ml = np.exp(-((x - 1500) / 200) ** 2) * 50 + np.exp(-((x - 3000) / 150) ** 2) * 80
    y_dft = np.exp(-((x - 1480) / 200) ** 2) * 48 + np.exp(-((x - 3050) / 150) ** 2) * 75

    ax.plot(x, y_ml, label='ML (MACE-MP + Espaloma)', linewidth=2, color='#2196F3')
    ax.plot(x, y_dft, label='DFT (B3LYP)', linewidth=2, linestyle='--', color='#F44336')

    ax.set_xlabel('Wavenumber (cm⁻¹)', fontsize=12)
    ax.set_ylabel('Intensity (km/mol)', fontsize=12)
    ax.set_title('IR Spectrum Comparison', fontsize=14, fontweight='bold')
    ax.legend(fontsize=11)
    ax.grid(alpha=0.3)

    plt.tight_layout()
    output_path = "../assets/plots/sample_spectrum.png"
    plt.savefig(output_path, dpi=150, bbox_inches='tight')
    plt.close()

    return output_path


def main():
    """Create a test PowerPoint presentation."""
    print("Creating test PowerPoint presentation...")
    print("=" * 60)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    print("Creating title slide...")
    create_title_slide(
        prs,
        "MACE-Gaussian Interface",
        "ML-Accelerated Anharmonic IR Spectroscopy\n\nYour Name\nResearch Group Meeting\nNovember 2025"
    )

    # Slide 2: Motivation
    print("Creating motivation slide...")
    create_content_slide(
        prs,
        "Motivation & Problem",
        [
            (0, "Anharmonic IR spectroscopy is essential for accurate predictions"),
            (1, "Includes overtones, combination bands, Fermi resonance"),
            (0, "DFT anharmonic calculations are computationally expensive"),
            (1, "Hours to days per molecule"),
            (1, "Prohibitive for high-throughput screening"),
            (0, "Can we use ML to accelerate without sacrificing accuracy?"),
        ]
    )

    # Slide 3: Solution Overview
    print("Creating solution overview slide...")
    create_content_slide(
        prs,
        "Our Solution: Hybrid ML-QM Approach",
        [
            (0, "Use MACE ML potentials for fast calculations"),
            (1, "Geometry optimization"),
            (1, "Force constants (Hessian)"),
            (1, "Dipole moment derivatives"),
            (0, "Inject ML data into Gaussian for anharmonic corrections"),
            (1, "Perturbation theory (VPT2)"),
            (1, "Overtones and combination bands"),
            (0, "Result: 10-100× speedup with comparable accuracy"),
        ]
    )

    # Slide 4: Architecture diagram (placeholder)
    print("Creating architecture slide...")
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    add_styled_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
        "System Architecture",
        font_size=32, bold=True
    )
    add_styled_textbox(
        slide, Inches(2), Inches(2.5), Inches(6), Inches(2),
        "[Flowchart Placeholder]\n\nXYZ → MACE Opt → Gaussian Freq → ZMQ Bridge → ML Dipoles → Analysis → Report",
        font_size=16, color=RGBColor(100, 100, 100)
    )

    # Slide 5: ZMQ Bridge (two column)
    print("Creating ZMQ bridge slide...")
    create_two_column_slide(
        prs,
        "Technical Innovation: ZMQ Bridge",
        left_content=[
            (0, "The Challenge"),
            (1, "Gaussian (Fortran) needs Python ML"),
            (1, "No direct FFI"),
            (0, "The Solution"),
            (1, "Inter-Process Communication"),
            (1, "ZMQ (ZeroMQ) message passing"),
            (1, "Unix socket (.ipc file)"),
        ],
        right_content=[
            (0, "How It Works"),
            (1, "Main script starts ZMQ server"),
            (1, "Gaussian calls helper script"),
            (1, "Helper connects to ZMQ socket"),
            (1, "Python calculates dipoles"),
            (1, "Data injected back to Gaussian"),
            (0, "Real-time bidirectional communication!"),
        ]
    )

    # Slide 6: Code example (placeholder)
    print("Creating code example slide...")
    slide = create_code_slide(
        prs,
        "ZMQ Bridge Implementation",
        "../assets/code_snippets/vscode_dark_zmq.png"
    )

    # Slide 7: Module swapping
    print("Creating module swapping slide...")
    create_content_slide(
        prs,
        "Technical Challenge: Dual MACE Packages",
        [
            (0, "Need two incompatible MACE versions:"),
            (1, "Standard MACE: energy/forces (optimization)"),
            (1, "Custom MACE: dipole moments (spectroscopy)"),
            (0, "Solution: Dynamic module swapping"),
            (1, "Monkey-patch sys.modules"),
            (1, "Swap implementations on-the-fly"),
            (1, "fake_module_from_real() mechanism"),
            (0, "Enables seamless switching between calculators"),
        ]
    )

    # Slide 8: Sample plot
    print("Creating sample plot...")
    plot_path = create_sample_plot()
    create_image_slide(
        prs,
        "Results: Spectrum Comparison",
        plot_path,
        "Water molecule: ML vs DFT anharmonic IR spectra"
    )

    # Slide 9: Validation
    print("Creating validation slide...")
    create_content_slide(
        prs,
        "Validation & Analysis",
        [
            (0, "Rigorous mode-by-mode comparison"),
            (1, "Eigenvector-based mode matching (not frequency)"),
            (1, "Handles mode crossing and frequency errors"),
            (0, "Statistical metrics"),
            (1, "MAE, RMSE, R² for frequencies"),
            (1, "Intensity correlations"),
            (1, "Mode overlap heatmaps"),
            (0, "Multiple ML methods tested"),
            (1, "MACE-MP, MACE-OMOL × Espaloma, MACE-ML dipoles"),
        ]
    )

    # Slide 10: Impact & Future
    print("Creating impact slide...")
    create_content_slide(
        prs,
        "Impact & Future Directions",
        [
            (0, "Current Achievements"),
            (1, "10-100× speedup over DFT"),
            (1, "Comparable accuracy for small molecules"),
            (1, "Automated workflow with HTML reports"),
            (0, "Future Work"),
            (1, "Extend to larger molecules (drugs, polymers)"),
            (1, "Improve dipole ML models"),
            (1, "High-throughput spectral database"),
            (1, "Integration with experimental workflows"),
        ]
    )

    # Slide 11: Thank you
    print("Creating closing slide...")
    create_section_header(prs, "Thank You!\n\nQuestions?")

    # Save presentation
    output_file = "../test_presentation.pptx"
    prs.save(output_file)

    print("\n" + "=" * 60)
    print(f"Test PowerPoint created: {output_file}")
    print(f"Total slides: {len(prs.slides)}")
    print("\nNow run the test scripts to generate assets:")
    print("  python test_code_styles.py")
    print("  python test_flowcharts.py")


if __name__ == "__main__":
    main()
