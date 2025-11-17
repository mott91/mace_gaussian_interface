#!/usr/bin/env python3
"""
Generate technical detail slides for MACE-Gaussian Interface presentation.

Covers:
1. Dipole calculation methods (MACE-ML and Espaloma)
2. Eigenvector mode matching algorithm
3. Output file parsing requirements
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Terminal Dark Color Scheme (GitHub Dark)
BG = RGBColor(13, 17, 23)        # #0D1117 - background
TEXT = RGBColor(201, 209, 217)   # #C9D1D9 - main text
ACCENT = RGBColor(88, 166, 255)  # #58A6FF - blue accent
DIM = RGBColor(139, 148, 158)    # #8B949E - dimmed text
GREEN = RGBColor(87, 171, 90)    # #57AB5A - success/input
ORANGE = RGBColor(219, 109, 40)  # #DB6D28 - warning/highlight
PURPLE = RGBColor(188, 140, 255) # #BC8CFF - special

FONT = "Consolas"


def create_dipole_calculation_slide(prs):
    """Slide explaining dipole calculation methods."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat dipole_calculation_methods.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(6.0))
    tf = content_box.text_frame
    tf.word_wrap = False

    lines = [
        ("DIPOLE MOMENT CALCULATION METHODS", ACCENT),
        ("═════════════════════════════════════════════════════════════════════════", DIM),
        ("", TEXT),
        ("1. MACE-ML (Custom Dipole-Enabled MACE)", ORANGE),
        ("   ─────────────────────────────────────", DIM),
        ("   Model: DipolePolarizabilityMACE", TEXT),
        ("   Path:  ~/mace_gaussian/dipole_model/model_1.model", DIM),
        ("", TEXT),
        ("   Architecture:", TEXT),
        ("   ┌──────────────────────────────────────────────────────────────┐", ACCENT),
        ("   │  Input: Atomic coordinates (x, y, z) + atomic numbers       │", TEXT),
        ("   │         ↓                                                    │", DIM),
        ("   │  E(3)-equivariant message passing layers                    │", TEXT),
        ("   │         ↓                                                    │", DIM),
        ("   │  Output: μ = (μx, μy, μz)  [Debye]                          │", GREEN),
        ("   │          Direct prediction of molecular dipole moment        │", TEXT),
        ("   └──────────────────────────────────────────────────────────────┘", ACCENT),
        ("", TEXT),
        ("   Advantages: Fast (~0.1s per geometry), physically consistent", TEXT),
        ("   Limitations: Requires custom-trained model", DIM),
        ("", TEXT),
        ("", TEXT),
        ("2. ESPALOMA (Charge-Based Approach)", ORANGE),
        ("   ──────────────────────────────────", DIM),
        ("   Method: Atomic partial charges from graph neural network", TEXT),
        ("", TEXT),
        ("   Algorithm:", TEXT),
        ("   ┌──────────────────────────────────────────────────────────────┐", ACCENT),
        ("   │  1. Predict partial charges: q₁, q₂, ..., qₙ               │", TEXT),
        ("   │                                                              │", TEXT),
        ("   │  2. Calculate dipole moment:                                │", TEXT),
        ("   │                        n                                     │", TEXT),
        ("   │     μ = Σ qᵢ · rᵢ     (vector sum)                          │", GREEN),
        ("   │        i=1                                                   │", TEXT),
        ("   │                                                              │", TEXT),
        ("   │     where: qᵢ = partial charge on atom i                    │", DIM),
        ("   │            rᵢ = position vector of atom i                   │", DIM),
        ("   └──────────────────────────────────────────────────────────────┘", ACCENT),
        ("", TEXT),
        ("   Advantages: Widely available, interpretable (charge-based)", TEXT),
        ("   Limitations: Indirect approach, charge assignment ambiguity", DIM),
        ("", TEXT),
        ("", TEXT),
        ("FALLBACK HIERARCHY: MACE-ML → Espaloma → xTB → Geometry-based", PURPLE),
    ]

    for i, (line, color) in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(9)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(0)
        p.line_spacing = 1.0


def create_eigenvector_matching_slide(prs):
    """Slide explaining eigenvector mode matching."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat eigenvector_mode_matching.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(6.0))
    tf = content_box.text_frame
    tf.word_wrap = False

    lines = [
        ("EIGENVECTOR MODE MATCHING ALGORITHM", ACCENT),
        ("═════════════════════════════════════════════════════════════════════════", DIM),
        ("", TEXT),
        ("PROBLEM: Frequency Ordering Mismatch", ORANGE),
        ("─────────────────────────────────────", DIM),
        ("", TEXT),
        ("DFT and ML calculations produce vibrational frequencies in different orders:", TEXT),
        ("", TEXT),
        ("  DFT Frequencies:  [500, 1200, 1500, 2900, 3100, 3200] cm⁻¹", GREEN),
        ("  ML Frequencies:   [498, 3195, 1205, 2895, 1498, 3098] cm⁻¹", ORANGE),
        ("                     ↑    ↑     ↑     ↑     ↑     ↑", DIM),
        ("  Cannot directly compare by index!", DIM),
        ("", TEXT),
        ("", TEXT),
        ("SOLUTION: Eigenvector Dot Product Matching", ACCENT),
        ("───────────────────────────────────────────", DIM),
        ("", TEXT),
        ("Each vibrational mode has an eigenvector describing atomic displacements:", TEXT),
        ("", TEXT),
        ("┌────────────────────────────────────────────────────────────────────────┐", ACCENT),
        ("│  For each ML mode i:                                                   │", TEXT),
        ("│                                                                         │", TEXT),
        ("│    1. Calculate dot product with all DFT modes:                        │", TEXT),
        ("│                                                                         │", TEXT),
        ("│       similarity[j] = |v_ML[i] · v_DFT[j]|                             │", GREEN),
        ("│                                                                         │", TEXT),
        ("│       where v_ML[i], v_DFT[j] are normalized eigenvectors              │", DIM),
        ("│                                                                         │", TEXT),
        ("│    2. Find best match:                                                 │", TEXT),
        ("│                                                                         │", TEXT),
        ("│       matched_mode[i] = argmax similarity[j]                           │", GREEN),
        ("│                              j                                          │", TEXT),
        ("│                                                                         │", TEXT),
        ("│    3. If similarity > threshold (0.8):                                 │", TEXT),
        ("│       → Match confirmed, compare frequencies                           │", TEXT),
        ("│       Else:                                                            │", TEXT),
        ("│       → Mode not found in DFT calculation                              │", TEXT),
        ("└────────────────────────────────────────────────────────────────────────┘", ACCENT),
        ("", TEXT),
        ("", TEXT),
        ("EXAMPLE MATCHING:", ORANGE),
        ("", TEXT),
        ("  ML Mode 1 (498 cm⁻¹)  → DFT Mode 1 (500 cm⁻¹)   [dot product: 0.996]", TEXT),
        ("  ML Mode 2 (3195 cm⁻¹) → DFT Mode 6 (3200 cm⁻¹)  [dot product: 0.989]", TEXT),
        ("  ML Mode 3 (1205 cm⁻¹) → DFT Mode 2 (1200 cm⁻¹)  [dot product: 0.993]", TEXT),
        ("  ...", DIM),
        ("", TEXT),
        ("Implementation: gaussian_parser.py, analyze_spectra.py", PURPLE),
    ]

    for i, (line, color) in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(9)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(0)
        p.line_spacing = 1.0


def create_parsing_slide(prs):
    """Slide explaining output file parsing requirements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat output_file_parsing.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(6.0))
    tf = content_box.text_frame
    tf.word_wrap = False

    lines = [
        ("OUTPUT FILE PARSING REQUIREMENTS", ACCENT),
        ("═════════════════════════════════════════════════════════════════════════", DIM),
        ("", TEXT),
        ("MULTIPLE FILE FORMATS TO PARSE:", ORANGE),
        ("", TEXT),
        ("1. GAUSSIAN OUTPUT FILES (.log)", ACCENT),
        ("   ─────────────────────────────", DIM),
        ("   Path: comparison_results/{molecule}/wb97xd/{molecule}_freq_anharm.log", TEXT),
        ("   Size: ~2-10 MB (text), 10,000-50,000 lines", DIM),
        ("", TEXT),
        ("   Parsing Targets:", TEXT),
        ("   ┌──────────────────────────────────────────────────────────────────┐", ACCENT),
        ("   │ • Harmonic frequencies:                                          │", TEXT),
        ("   │   Pattern: 'Frequencies --  500.00  1200.00  1500.00'           │", DIM),
        ("   │                                                                  │", TEXT),
        ("   │ • Anharmonic frequencies:                                        │", TEXT),
        ("   │   Pattern: 'Fundamental Bands (cm-1)'                            │", DIM),
        ("   │                                                                  │", TEXT),
        ("   │ • Overtones and combination bands:                               │", TEXT),
        ("   │   Pattern: '2(1)' or '1(1) + 2(1)'                               │", DIM),
        ("   │                                                                  │", TEXT),
        ("   │ • IR intensities:                                                │", TEXT),
        ("   │   Pattern: 'IR Inten    --   25.00   50.00  100.00'             │", DIM),
        ("   │                                                                  │", TEXT),
        ("   │ • Normal mode eigenvectors:                                      │", TEXT),
        ("   │   3N-6 modes × N atoms × 3 coordinates                           │", DIM),
        ("   │   Pattern: 'Atom  AN      X      Y      Z'                       │", DIM),
        ("   └──────────────────────────────────────────────────────────────────┘", ACCENT),
        ("", TEXT),
        ("   Challenges:", ORANGE),
        ("   • Multi-column format (3 modes per block)", TEXT),
        ("   • Scattered data across file (search 1000s of lines)", TEXT),
        ("   • Anharmonic section sometimes missing (wb97xd issues)", TEXT),
        ("", TEXT),
        ("", TEXT),
        ("2. JSON RESULTS FILES (ML calculations)", ACCENT),
        ("   ──────────────────────────────────────", DIM),
        ("   Paths: comparison_results/{molecule}/{energy}_{dipole}/results.json", TEXT),
        ("   Examples:", DIM),
        ("     • mace_mp_espaloma/results.json", DIM),
        ("     • mace_omol_mace_ml/results.json", DIM),
        ("", TEXT),
        ("   Structure:", TEXT),
        ("   ┌──────────────────────────────────────────────────────────────────┐", ACCENT),
        ("   │ {                                                                │", TEXT),
        ("   │   \"calculator_type\": \"ml\",                                      │", DIM),
        ("   │   \"energy_calculator\": \"mace_mp\",                               │", DIM),
        ("   │   \"dipole_calculator\": \"espaloma\",                              │", DIM),
        ("   │   \"frequencies\": {                                               │", DIM),
        ("   │     \"harmonic\": [500.0, 1200.0, 1500.0, ...],                   │", GREEN),
        ("   │     \"anharmonic\": [498.0, 1195.0, 1497.0, ...]                  │", GREEN),
        ("   │   },                                                             │", DIM),
        ("   │   \"ir_intensities\": {                                            │", DIM),
        ("   │     \"harmonic\": [25.0, 50.0, 100.0, ...],                       │", GREEN),
        ("   │     \"anharmonic\": [26.0, 51.0, 102.0, ...]                      │", GREEN),
        ("   │   },                                                             │", DIM),
        ("   │   \"eigenvectors\": [[...], [...], ...],                          │", DIM),
        ("   │   \"runtime_s\": 45.2                                             │", DIM),
        ("   │ }                                                                │", TEXT),
        ("   └──────────────────────────────────────────────────────────────────┘", ACCENT),
        ("", TEXT),
        ("", TEXT),
        ("PARSER IMPLEMENTATION:", PURPLE),
        ("  • gaussian_parser.py: Regex-based Gaussian log parsing", TEXT),
        ("  • results_manager.py: JSON loading and validation", TEXT),
        ("  • comparison_workflow.py: Aggregates all results for analysis", TEXT),
    ]

    for i, (line, color) in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(8)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(0)
        p.line_spacing = 1.0


def main():
    """Generate technical details presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Generating technical detail slides...")

    create_dipole_calculation_slide(prs)
    print("  ✓ Dipole calculation methods (MACE-ML & Espaloma)")

    create_eigenvector_matching_slide(prs)
    print("  ✓ Eigenvector mode matching algorithm")

    create_parsing_slide(prs)
    print("  ✓ Output file parsing requirements")

    # Save presentation
    output_path = "../technical_details.pptx"
    prs.save(output_path)
    print(f"\n✓ Technical details presentation created: {output_path}")
    print(f"  → 3 slides covering implementation details")
    print(f"  → Terminal Dark styling")


if __name__ == "__main__":
    main()
