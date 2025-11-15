#!/usr/bin/env python3
"""
Generate multiple minimal presentation templates.
Each template has 2-3 slides showcasing a different aesthetic.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ============================================================================
# TEMPLATE 1: TERMINAL DARK
# ============================================================================

def create_terminal_dark_template():
    """Dark terminal theme - hacker aesthetic."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    BG = RGBColor(13, 17, 23)          # #0D1117 GitHub dark
    TEXT = RGBColor(201, 209, 217)     # #C9D1D9
    ACCENT = RGBColor(88, 166, 255)    # #58A6FF bright blue
    DIM = RGBColor(139, 148, 158)      # #8B949E

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Prompt-style title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = "$ ./presentation.sh\n> MACE-Gaussian Interface"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = 'Consolas'
    p.font.color.rgb = ACCENT
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.6))
    tf = subtitle_box.text_frame
    tf.text = "// ML-accelerated anharmonic spectroscopy"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.name = 'Consolas'
    p.font.color.rgb = DIM

    # Info
    info_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.6))
    tf = info_box.text_frame
    tf.text = "Your Name · Research Group · 2025-11-15"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.name = 'Consolas'
    p.font.color.rgb = DIM
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Content
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title with prompt
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat motivation.md"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.name = 'Consolas'
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
    tf = content_box.text_frame

    lines = [
        ("# Problem", ACCENT),
        ("  → DFT anharmonic: computationally expensive", TEXT),
        ("  → Hours to days per molecule", TEXT),
        ("", TEXT),
        ("# Solution", ACCENT),
        ("  → ML potentials: 10-100× faster", TEXT),
        ("  → Maintain accuracy via hybrid approach", TEXT),
        ("", TEXT),
        ("# Impact", ACCENT),
        ("  → High-throughput spectral predictions", TEXT),
    ]

    for i, (line, color) in enumerate(lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(14)
        p.font.name = 'Consolas'
        p.font.color.rgb = color
        p.space_after = Pt(6)

    # Slide 3: Image placeholder
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ ls -la results/"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.name = 'Consolas'
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Image placeholder
    rect = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(7), Inches(5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(22, 27, 34)
    rect.line.color.rgb = RGBColor(48, 54, 61)

    placeholder = slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(4), Inches(1))
    tf = placeholder.text_frame
    tf.text = "[flowchart / plot goes here]"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.name = 'Consolas'
    p.font.color.rgb = DIM
    p.alignment = PP_ALIGN.CENTER

    prs.save("../template_terminal_dark.pptx")
    print("✓ Terminal Dark template created (3 slides)")


# ============================================================================
# TEMPLATE 2: MONOCHROME MINIMAL
# ============================================================================

def create_monochrome_template():
    """Pure black & white minimal template."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(128, 128, 128)
    LIGHT_GRAY = RGBColor(240, 240, 240)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Large title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = title_box.text_frame
    tf.text = "MACE-Gaussian\nInterface"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.name = 'Helvetica'
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 0.9

    # Thin line
    line = slide.shapes.add_connector(1, Inches(1), Inches(4.3), Inches(4), Inches(4.3))
    line.line.color.rgb = BLACK
    line.line.width = Pt(2)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.6), Inches(8), Inches(1))
    tf = subtitle_box.text_frame
    tf.text = "ML-Accelerated Anharmonic IR Spectroscopy"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.name = 'Helvetica'
    p.font.color.rgb = GRAY

    # Slide 2: Content
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.6))
    tf = title_box.text_frame
    tf.text = "Overview"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.name = 'Helvetica'
    p.font.color.rgb = BLACK
    p.font.bold = True

    # Line under title
    line = slide.shapes.add_connector(1, Inches(0.7), Inches(1.2), Inches(9.3), Inches(1.2))
    line.line.color.rgb = BLACK
    line.line.width = Pt(1)

    # Minimalist bullets
    y_pos = 1.8
    bullets = [
        "Hybrid ML-QM approach",
        "10-100× computational speedup",
        "Comparable accuracy to pure DFT",
        "Novel ZMQ bridge architecture",
        "Automated validation pipeline",
    ]

    for bullet in bullets:
        # Dot
        dot_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(0.2), Inches(0.3))
        tf = dot_box.text_frame
        tf.text = "·"
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK

        # Text
        text_box = slide.shapes.add_textbox(Inches(1.4), Inches(y_pos), Inches(7.5), Inches(0.4))
        tf = text_box.text_frame
        tf.text = bullet
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.name = 'Helvetica'
        p.font.color.rgb = BLACK

        y_pos += 0.65

    # Slide 3: Large number slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Giant number
    num_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(3))
    tf = num_box.text_frame
    tf.text = "100×"
    p = tf.paragraphs[0]
    p.font.size = Pt(120)
    p.font.name = 'Helvetica'
    p.font.color.rgb = BLACK
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Description
    desc_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.5), Inches(3.5), Inches(2))
    tf = desc_box.text_frame
    tf.text = "Faster than\ntraditional DFT\nanharmonic\ncalculations"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.name = 'Helvetica'
    p.font.color.rgb = GRAY
    p.line_spacing = 1.3

    prs.save("../template_monochrome.pptx")
    print("✓ Monochrome Minimal template created (3 slides)")


# ============================================================================
# TEMPLATE 3: SCIENTIFIC PAPER STYLE
# ============================================================================

def create_scientific_template():
    """LaTeX/academic paper inspired template."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    TEXT = RGBColor(33, 33, 33)
    GRAY = RGBColor(100, 100, 100)
    LIGHT_GRAY = RGBColor(220, 220, 220)
    ACCENT = RGBColor(0, 51, 102)  # Dark blue

    # Slide 1: Title (like a paper title page)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "MACE-Gaussian Interface:\nML-Accelerated Anharmonic IR Spectroscopy"
    p = tf.paragraphs[0]
    p.font.size = Pt(26)
    p.font.name = 'Times New Roman'
    p.font.color.rgb = TEXT
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.2

    # Author info
    author_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
    tf = author_box.text_frame
    tf.text = "Your Name\nResearch Group, Institution\n\nNovember 2025"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.name = 'Times New Roman'
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.4

    # Slide 2: Abstract style
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section header
    header_box = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(8), Inches(0.5))
    tf = header_box.text_frame
    tf.text = "1. Introduction"
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.name = 'Times New Roman'
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Body text (justified, like a paper)
    body_box = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(8), Inches(5))
    tf = body_box.text_frame
    tf.word_wrap = True

    paragraphs_text = [
        "Anharmonic vibrational spectroscopy calculations using density functional theory (DFT) "
        "provide accurate predictions but remain computationally prohibitive for high-throughput "
        "applications, often requiring hours to days per molecule.",

        "We present a hybrid machine learning-quantum mechanics (ML-QM) approach that accelerates "
        "anharmonic IR spectroscopy by 10-100× while maintaining comparable accuracy to pure DFT methods.",

        "Key contributions:",
    ]

    for i, text in enumerate(paragraphs_text):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = text
        p.font.size = Pt(14)
        p.font.name = 'Times New Roman'
        p.font.color.rgb = TEXT
        p.alignment = PP_ALIGN.JUSTIFY
        p.space_after = Pt(12)

    # Numbered list
    items = [
        "Novel ZMQ-based inter-process communication bridge",
        "Dynamic module swapping for dual MACE implementations",
        "Eigenvector-based mode matching for validation",
    ]

    y_pos = 4.2
    for i, item in enumerate(items):
        # Number
        num_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(0.3), Inches(0.3))
        tf = num_box.text_frame
        tf.text = f"{i+1}."
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.name = 'Times New Roman'
        p.font.color.rgb = TEXT

        # Text
        text_box = slide.shapes.add_textbox(Inches(1.9), Inches(y_pos), Inches(7), Inches(0.4))
        tf = text_box.text_frame
        tf.text = item
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.name = 'Times New Roman'
        p.font.color.rgb = TEXT

        y_pos += 0.45

    prs.save("../template_scientific.pptx")
    print("✓ Scientific Paper template created (2 slides)")


# ============================================================================
# TEMPLATE 4: BRUTALIST/CODE AESTHETIC
# ============================================================================

def create_brutalist_template():
    """Brutalist design - bold, geometric, high contrast."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    RED = RGBColor(255, 0, 0)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK

    # Large block letters
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    tf = title_box.text_frame
    tf.text = "MACE\nGAUSSIAN\nINTERFACE"
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.name = 'Arial Black'
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 0.85

    # Red accent line
    rect = slide.shapes.add_shape(1, Inches(0.5), Inches(4.8), Inches(3), Inches(0.1))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RED
    rect.line.fill.background()

    # Small text
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(5), Inches(0.5))
    tf = subtitle_box.text_frame
    tf.text = "ML / QM HYBRID SYSTEM"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.name = 'Arial'
    p.font.color.rgb = WHITE
    p.font.bold = True

    # Slide 2: Content
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Split screen - black left, white right
    left_rect = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(5), Inches(7.5))
    left_rect.fill.solid()
    left_rect.fill.fore_color.rgb = BLACK
    left_rect.line.fill.background()

    # Left side - white text
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(4), Inches(1))
    tf = left_title.text_frame
    tf.text = "PROBLEM"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.name = 'Arial Black'
    p.font.color.rgb = WHITE
    p.font.bold = True

    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4))
    tf = left_content.text_frame
    tf.text = "DFT SLOW\nEXPENSIVE\nLIMITED SCALE"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = WHITE
    p.line_spacing = 1.8

    # Right side - black text
    right_title = slide.shapes.add_textbox(Inches(5.5), Inches(0.8), Inches(4), Inches(1))
    tf = right_title.text_frame
    tf.text = "SOLUTION"
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.name = 'Arial Black'
    p.font.color.rgb = BLACK
    p.font.bold = True

    right_content = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
    tf = right_content.text_frame
    tf.text = "ML FAST\nHYBRID ACCURATE\n100× SPEEDUP"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.name = 'Arial'
    p.font.color.rgb = BLACK
    p.line_spacing = 1.8

    # Red accent
    rect = slide.shapes.add_shape(1, Inches(4.95), Inches(0), Inches(0.1), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RED
    rect.line.fill.background()

    prs.save("../template_brutalist.pptx")
    print("✓ Brutalist template created (2 slides)")


def main():
    """Generate all templates."""
    print("Generating presentation templates...")
    print("=" * 60)

    create_terminal_dark_template()
    create_monochrome_template()
    create_scientific_template()
    create_brutalist_template()

    print("\n" + "=" * 60)
    print("All templates created!")
    print("\nTemplates:")
    print("  1. template_terminal_dark.pptx   - Dark terminal/hacker theme")
    print("  2. template_monochrome.pptx      - Pure B&W minimal")
    print("  3. template_scientific.pptx      - Academic paper style")
    print("  4. template_brutalist.pptx       - Bold geometric design")


if __name__ == "__main__":
    main()
