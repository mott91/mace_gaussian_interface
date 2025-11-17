#!/usr/bin/env python3
"""
Test different ASCII diagrams for the ZMQ bridge architecture.
Show the non-linear bidirectional communication flow.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# Terminal Dark color scheme
BG = RGBColor(13, 17, 23)          # #0D1117 GitHub dark
TEXT = RGBColor(201, 209, 217)     # #C9D1D9
ACCENT = RGBColor(88, 166, 255)    # #58A6FF bright blue
DIM = RGBColor(139, 148, 158)      # #8B949E
GREEN = RGBColor(87, 171, 90)      # #57AB5A
ORANGE = RGBColor(219, 109, 40)    # #DB6D28
FONT = 'Consolas'


def create_circular_zmq_diagram(prs):
    """Circular diagram showing bidirectional communication."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat zmq_bridge_circular.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Diagram box
    diagram_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(5.5))
    tf = diagram_box.text_frame

    diagram_lines = [
        ("┌─────────────────────────────────────────────┐", ACCENT),
        ("│  Main Python Process (gm_main.py)          │", ACCENT),
        ("│  • Creates ZMQ socket (server)             │", TEXT),
        ("│  • Launches Gaussian subprocess            │", TEXT),
        ("│  • Calculates ML dipoles                   │", TEXT),
        ("└─────────────────────────────────────────────┘", ACCENT),
        ("         │                        ▲", DIM),
        ("         │ launches               │ ZMQ", DIM),
        ("         ▼                        │ socket", DIM),
        ("┌─────────────────────────────────────────────┐", GREEN),
        ("│  Gaussian 16 (Fortran)                      │", GREEN),
        ("│  • Frequency calculation                    │", TEXT),
        ("│  • Calls External=\"gm_helper.py\"           │", TEXT),
        ("└─────────────────────────────────────────────┘", GREEN),
        ("         │                        ▲", DIM),
        ("         │ calls                  │ returns", DIM),
        ("         ▼                        │", DIM),
        ("┌─────────────────────────────────────────────┐", ORANGE),
        ("│  Helper Script (gm_helper.py)               │", ORANGE),
        ("│  • Connects to ZMQ socket                   │", TEXT),
        ("│  • Sends geometry file paths                │", TEXT),
        ("│  • Waits for 'done' signal                  │", TEXT),
        ("│  • Returns to Gaussian                      │", TEXT),
        ("└─────────────────────────────────────────────┘", ORANGE),
    ]

    for i, (line, color) in enumerate(diagram_lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(11)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(2)
        p.line_spacing = 1.0


def create_side_by_side_zmq_diagram(prs):
    """Side-by-side boxes with arrows showing communication."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat zmq_bridge_side_by_side.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Diagram box
    diagram_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = diagram_box.text_frame

    diagram_lines = [
        ("┌─────────────────────┐        ┌─────────────────────┐", DIM),
        ("│  Main Process       │◀──────▶│  Gaussian 16        │", DIM),
        ("│  (Python)           │  ZMQ   │  (Fortran)          │", DIM),
        ("│                     │        │                     │", DIM),
        ("│  • ZMQ server       │        │  • Freq calculation │", DIM),
        ("│  • ML dipole calc   │        │  • Calls external   │", DIM),
        ("└─────────────────────┘        └─────────────────────┘", DIM),
        ("         ▲                                 │", DIM),
        ("         │                                 │", DIM),
        ("         │         ZMQ socket              │", DIM),
        ("         │                                 ▼", DIM),
        ("         │          ┌─────────────────────────────┐", DIM),
        ("         └──────────│  Helper (gm_helper.py)      │", DIM),
        ("                    │  • Bridge between processes │", DIM),
        ("                    │  • Sends file paths         │", DIM),
        ("                    │  • Waits for completion     │", DIM),
        ("                    └─────────────────────────────┘", DIM),
        ("", TEXT),
        ("Flow:", ACCENT),
        ("  1. Main launches Gaussian with External directive", TEXT),
        ("  2. Gaussian calls helper when dipoles needed", TEXT),
        ("  3. Helper connects to main via ZMQ", TEXT),
        ("  4. Main calculates dipoles, writes output", TEXT),
        ("  5. Helper returns to Gaussian", TEXT),
    ]

    for i, (line, color) in enumerate(diagram_lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(11)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(2)


def create_detailed_zmq_diagram(prs):
    """Detailed diagram with numbered steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat zmq_bridge_detailed.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Diagram box
    diagram_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(8.4), Inches(5.9))
    tf = diagram_box.text_frame

    diagram_lines = [
        ("╔═══════════════════════════════════════════════════════╗", ACCENT),
        ("║  Main Process: gm_main.py (Python + ASE + MACE)      ║", ACCENT),
        ("╚═══════════════════════════════════════════════════════╝", ACCENT),
        ("         │ (1) Create ZMQ socket, bind to .ipc_file", DIM),
        ("         │ (2) Generate Gaussian input", DIM),
        ("         ▼", DIM),
        ("╔═══════════════════════════════════════════════════════╗", GREEN),
        ("║  Gaussian 16: Frequency Calculation                   ║", GREEN),
        ("║  Route: # freq (anharm) external=\"gm_helper.py\"       ║", TEXT),
        ("╚═══════════════════════════════════════════════════════╝", GREEN),
        ("         │ (3) During freq calc: call external program", DIM),
        ("         ▼", DIM),
        ("╔═══════════════════════════════════════════════════════╗", ORANGE),
        ("║  Helper Script: gm_helper.py                          ║", ORANGE),
        ("╚═══════════════════════════════════════════════════════╝", ORANGE),
        ("         │ (4) Connect to ZMQ socket (.ipc_file)", DIM),
        ("         │ (5) Send: \"infile|outfile\"", DIM),
        ("         ▼", DIM),
        ("╔═══════════════════════════════════════════════════════╗", ACCENT),
        ("║  Main Process: ZMQ Handler                            ║", ACCENT),
        ("╚═══════════════════════════════════════════════════════╝", ACCENT),
        ("         │ (6) Read geometry from infile", DIM),
        ("         │ (7) Calculate ML dipoles", DIM),
        ("         │ (8) Write Gaussian format to outfile", DIM),
        ("         │ (9) Send: \"done\"", DIM),
        ("         ▼", DIM),
        ("╔═══════════════════════════════════════════════════════╗", ORANGE),
        ("║  Helper: gm_helper.py                                 ║", ORANGE),
        ("╚═══════════════════════════════════════════════════════╝", ORANGE),
        ("         │ (10) Receive \"done\", return to Gaussian", DIM),
        ("         ▼", DIM),
        ("╔═══════════════════════════════════════════════════════╗", GREEN),
        ("║  Gaussian 16: Continue calculation                    ║", GREEN),
        ("╚═══════════════════════════════════════════════════════╝", GREEN),
    ]

    for i, (line, color) in enumerate(diagram_lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(9)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(1)
        p.line_spacing = 1.0


def create_simple_zmq_diagram(prs):
    """Simple single-line box diagram showing communication."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    tf.text = "$ cat zmq_bridge.txt"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    # Diagram box
    diagram_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(7.6), Inches(5))
    tf = diagram_box.text_frame

    diagram_lines = [
        ("┌───────────────────────────────────────┐", ACCENT),
        ("│  Main Process (gm_main.py)           │", ACCENT),
        ("│  → ZMQ server                         │", TEXT),
        ("│  → ML dipole calculations            │", TEXT),
        ("└───────────────────────────────────────┘", ACCENT),
        ("    ▲                              │", DIM),
        ("    │                              │ launches", DIM),
        ("    │ ZMQ socket                   ▼", DIM),
        ("    │         ┌───────────────────────────────────────┐", GREEN),
        ("    │         │  Gaussian 16                          │", GREEN),
        ("    │         │  → Frequency calculation              │", TEXT),
        ("    │         └───────────────────────────────────────┘", GREEN),
        ("    │                              │", DIM),
        ("    │                              │ calls external", DIM),
        ("    │                              ▼", DIM),
        ("    │         ┌───────────────────────────────────────┐", ORANGE),
        ("    └─────────│  Helper Script (gm_helper.py)         │", ORANGE),
        ("              │  → Connects main ↔ Gaussian           │", TEXT),
        ("              └───────────────────────────────────────┘", ORANGE),
    ]

    for i, (line, color) in enumerate(diagram_lines):
        if i > 0:
            tf.add_paragraph()
        p = tf.paragraphs[i]
        p.text = line
        p.font.size = Pt(12)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(2)
        p.line_spacing = 1.0


def main():
    """Generate test presentation with different ZMQ diagram styles."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Generating ZMQ bridge diagram test slides...")

    create_circular_zmq_diagram(prs)
    print("  ✓ Circular diagram (3 boxes vertical)")

    create_side_by_side_zmq_diagram(prs)
    print("  ✓ Side-by-side diagram (with flow description)")

    create_detailed_zmq_diagram(prs)
    print("  ✓ Detailed diagram (numbered steps)")

    create_simple_zmq_diagram(prs)
    print("  ✓ Simple diagram (single-line boxes)")

    # Save presentation
    output_path = "../test_zmq_diagrams.pptx"
    prs.save(output_path)
    print(f"\n✓ Test presentation created: {output_path}")
    print(f"  → 4 slides with different ZMQ diagram styles")
    print(f"  → Shows non-linear bidirectional communication")
    print(f"  → Terminal Dark styling")
    print(f"\nReview the slides and pick your favorite!")


if __name__ == "__main__":
    main()
