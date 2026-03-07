#!/usr/bin/env python3
"""
Title slide banner workshop — 8 style variations.
Run: python3.8 presentation/banners/generate_banners.py
Output: presentation/banners/title_banners.pptx

Styles:
  1 — MACE in doom block letters (pyfiglet)
  2 — ASCII IR spectrum art
  3 — system boot / progress bars
  4 — v1 spirit, perfected
  5 — lean font (thin strokes)
  6 — lean font (filled █)
  7 — block font (thin strokes)
  8 — block font (filled █)
  9 — block + GAUSSIAN, centered, neofetch spirit
"""

import pyfiglet
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

BG     = RGBColor(13,  17,  23)
TEXT   = RGBColor(201, 209, 217)
ACCENT = RGBColor(88,  166, 255)
GREEN  = RGBColor(63,  185, 80)
YELLOW = RGBColor(210, 153, 34)
DIM    = RGBColor(139, 148, 158)
FONT   = "Consolas"
DATE   = "Manuel Ott · Hofer Lab · 2026-03-26"


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def footer(slide):
    box = slide.shapes.add_textbox(Inches(1), Inches(6.85), Inches(8), Inches(0.4))
    tf = box.text_frame
    tf.text = DATE
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.name = FONT
    p.font.color.rgb = DIM
    p.alignment = PP_ALIGN.CENTER


def label(slide, n, desc):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(8), Inches(0.3))
    tf = box.text_frame
    tf.text = f"// style {n} — {desc}"
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.name = FONT
    p.font.color.rgb = DIM


def block(slide, lines, x, y, w, h, size, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = False
    for i, (text, color) in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(0)
        p.alignment = align


# ─────────────────────────────────────────────────────────────────────────────
# Style 1: MACE in doom font block letters
# ─────────────────────────────────────────────────────────────────────────────

def style_1(prs):
    slide = blank(prs)
    label(slide, 1, "figlet doom — MACE block letters")

    raw = pyfiglet.figlet_format("MACE", font="doom")
    art_lines = [(ln, ACCENT) for ln in raw.splitlines() if ln.strip()]

    block(slide, art_lines, x=0.6, y=1.0, w=9.2, h=2.8, size=16,
          align=PP_ALIGN.CENTER)

    block(slide, [
        ("─" * 44,                           DIM),
        ("",                                  DIM),
        ("$ ./presentation.sh",               ACCENT),
        ("",                                  DIM),
        ("// ML-accelerated IR spectroscopy", TEXT),
        ("// MACE × Gaussian 16 · ZMQ · VPT2", DIM),
    ], x=0.8, y=4.1, w=8.4, h=2.0, size=16, align=PP_ALIGN.CENTER)

    footer(slide)


# ─────────────────────────────────────────────────────────────────────────────
# Style 2: ASCII IR spectrum art
# ─────────────────────────────────────────────────────────────────────────────

def style_2(prs):
    slide = blank(prs)
    label(slide, 2, "ASCII IR spectrum")

    # Hand-crafted IR spectrum — 52 columns, 8 rows
    # Peaks: ~700 (aromatic), ~1000-1200 (C-O), ~1680 (C=O), ~3000 (C-H), ~3300 (O-H)
    H = [0] * 52
    for i, h in [
        (3,4),(4,5),(5,5),(6,4),           # ~700  aromatic C-H bend
        (8,5),(9,7),(10,8),(11,7),(12,5),  # ~1050 C-O stretch
        (14,4),(15,6),(16,5),              # ~1200 C-O
        (18,4),(19,5),(20,4),              # ~1380 C-H deform
        (21,8),(22,8),(23,7),(24,5),       # ~1680 C=O stretch (tallest)
        (33,3),(34,4),(35,3),              # ~2900 C-H stretch
        (36,5),(37,6),(38,5),              # ~3050 C-H stretch
        (40,3),(41,5),(42,6),(43,4),       # ~3300 O-H stretch
    ]:
        H[i] = h

    rows = []
    for row in range(8, 0, -1):
        line = "  "
        for h in H:
            line += "|" if h >= row else " "
        rows.append((line, ACCENT if row >= 6 else (TEXT if row >= 3 else DIM)))

    rows.append(("  " + "─" * 52 + "→", DIM))
    rows.append(("   400    700   1200  1680       2900 3050 3300    cm⁻¹", DIM))

    block(slide, rows, x=0.4, y=0.8, w=9.2, h=4.4, size=12, align=PP_ALIGN.LEFT)

    block(slide, [
        ("$ mace-gaussian run aspirin.xyz  # harmonic benchmark", ACCENT),
        ("",                                                        DIM),
        ("  // ML-accelerated IR spectroscopy",                    DIM),
    ], x=0.6, y=5.5, w=9.0, h=1.0, size=14, align=PP_ALIGN.LEFT)

    footer(slide)


# ─────────────────────────────────────────────────────────────────────────────
# Style 3: System boot / loading bars
# ─────────────────────────────────────────────────────────────────────────────

def style_3(prs):
    slide = blank(prs)
    label(slide, 3, "system boot sequence")

    BAR = "=" * 24
    OK  = "[done]"

    lines = [
        ("  Initializing mace-gaussian v1.1",                 ACCENT),
        ("",                                                    DIM),
        (f"  [{BAR}]  Loading MACE-OMOL-0        {OK}",       GREEN),
        (f"  [{BAR}]  Loading MACE4IR dipole      {OK}",      GREEN),
        (f"  [{BAR}]  Starting ZMQ IPC server     {OK}",      GREEN),
        (f"  [{BAR}]  Gaussian 16 bridge ready    {OK}",      GREEN),
        ("",                                                    DIM),
        ("  ─" * 22,                                           DIM),
        ("",                                                    DIM),
        ("  $ mace-gaussian run molecule.xyz",                 ACCENT),
        ("",                                                    DIM),
        ("  // all systems go",                                DIM),
    ]

    block(slide, lines, x=0.8, y=1.5, w=8.8, h=4.5, size=16,
          align=PP_ALIGN.LEFT)

    footer(slide)


# ─────────────────────────────────────────────────────────────────────────────
# Style 4: v1 spirit — pure, centered, breathing room
# ─────────────────────────────────────────────────────────────────────────────

def style_4(prs):
    slide = blank(prs)
    label(slide, 4, "v1 spirit — minimal + centered")

    block(slide, [
        ("$ ./presentation.sh",               ACCENT),
        ("",                                   DIM),
        ("",                                   DIM),
        ("// ML-accelerated IR spectroscopy",  DIM),
    ], x=1.0, y=2.6, w=8.0, h=2.2, size=30, align=PP_ALIGN.CENTER)

    footer(slide)


def _figlet_fill(text, font):
    """Generate figlet art with all strokes replaced by █."""
    raw = pyfiglet.figlet_format(text, font=font)
    filled = ""
    for ch in raw:
        if ch in ("_", "/", "\\", "|"):
            filled += "█"
        else:
            filled += ch
    return filled


# ─────────────────────────────────────────────────────────────────────────────
# Style 5: lean font — thin strokes
# ─────────────────────────────────────────────────────────────────────────────

def style_5(prs):
    slide = blank(prs)
    label(slide, 5, "lean — thin strokes")

    raw = pyfiglet.figlet_format("MACE", font="lean")
    art_lines = [(ln, ACCENT) for ln in raw.splitlines() if ln.strip()]

    block(slide, art_lines, x=0.4, y=1.0, w=9.2, h=2.8, size=16,
          align=PP_ALIGN.CENTER)

    block(slide, [
        ("─" * 44,                           DIM),
        ("",                                  DIM),
        ("$ mace-gaussian run molecule.xyz",  ACCENT),
        ("",                                  DIM),
        ("// ML-accelerated IR spectroscopy", TEXT),
        ("// MACE × Gaussian 16 · ZMQ · VPT2", DIM),
    ], x=0.8, y=4.1, w=8.4, h=2.0, size=16, align=PP_ALIGN.CENTER)

    footer(slide)


# ─────────────────────────────────────────────────────────────────────────────
# Style 6: lean font — filled █
# ─────────────────────────────────────────────────────────────────────────────

def style_6(prs):
    slide = blank(prs)
    label(slide, 6, "lean — filled █")

    raw = _figlet_fill("MACE", "lean")
    art_lines = [(ln, ACCENT) for ln in raw.splitlines() if ln.strip()]

    block(slide, art_lines, x=0.4, y=1.0, w=9.2, h=2.8, size=16,
          align=PP_ALIGN.CENTER)

    block(slide, [
        ("─" * 44,                           DIM),
        ("",                                  DIM),
        ("$ mace-gaussian run molecule.xyz",  ACCENT),
        ("",                                  DIM),
        ("// ML-accelerated IR spectroscopy", TEXT),
        ("// MACE × Gaussian 16 · ZMQ · VPT2", DIM),
    ], x=0.8, y=4.1, w=8.4, h=2.0, size=16, align=PP_ALIGN.CENTER)

    footer(slide)


# ─────────────────────────────────────────────────────────────────────────────
# Style 7: block font — thin strokes
# ─────────────────────────────────────────────────────────────────────────────

def style_7(prs):
    slide = blank(prs)
    label(slide, 7, "block — thin strokes")

    raw = pyfiglet.figlet_format("MACE", font="block")
    art_lines = [(ln, ACCENT) for ln in raw.splitlines() if ln.strip()]

    block(slide, art_lines, x=0.6, y=1.0, w=9.2, h=2.8, size=16,
          align=PP_ALIGN.CENTER)

    block(slide, [
        ("─" * 44,                           DIM),
        ("",                                  DIM),
        ("$ mace-gaussian run molecule.xyz",  ACCENT),
        ("",                                  DIM),
        ("// ML-accelerated IR spectroscopy", TEXT),
        ("// MACE × Gaussian 16 · ZMQ · VPT2", DIM),
    ], x=0.8, y=4.1, w=8.4, h=2.0, size=16, align=PP_ALIGN.CENTER)

    footer(slide)


# ─────────────────────────────────────────────────────────────────────────────
# Style 8: block font — filled █
# ─────────────────────────────────────────────────────────────────────────────

def style_8(prs):
    slide = blank(prs)
    label(slide, 8, "block — filled █")

    raw = _figlet_fill("MACE", "block")
    art_lines = [(ln, ACCENT) for ln in raw.splitlines() if ln.strip()]

    block(slide, art_lines, x=0.6, y=1.0, w=9.2, h=2.8, size=16,
          align=PP_ALIGN.CENTER)

    block(slide, [
        ("─" * 44,                           DIM),
        ("",                                  DIM),
        ("$ mace-gaussian run molecule.xyz",  ACCENT),
        ("",                                  DIM),
        ("// ML-accelerated IR spectroscopy", TEXT),
        ("// MACE × Gaussian 16 · ZMQ · VPT2", DIM),
    ], x=0.8, y=4.1, w=8.4, h=2.0, size=16, align=PP_ALIGN.CENTER)

    footer(slide)


# ─────────────────────────────────────────────────────────────────────────────
# Style 9: block + GAUSSIAN subtitle, centered, neofetch spirit
# ─────────────────────────────────────────────────────────────────────────────

def style_9(prs):
    slide = blank(prs)
    label(slide, 9, "block + GAUSSIAN — neofetch spirit")

    raw = pyfiglet.figlet_format("MACE", font="block")
    art_lines = [(ln, ACCENT) for ln in raw.splitlines() if ln.strip()]
    art_lines.append(("       GAUSSIAN", ACCENT))

    block(slide, art_lines, x=0.6, y=1.6, w=9.2, h=3.0, size=16,
          align=PP_ALIGN.CENTER)

    block(slide, [
        ("// ML-accelerated IR spectrum calculations", DIM),
    ], x=0.8, y=5.0, w=8.4, h=0.6, size=16, align=PP_ALIGN.CENTER)

    footer(slide)


# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    style_1(prs)
    style_2(prs)
    style_3(prs)
    style_4(prs)
    style_5(prs)
    style_6(prs)
    style_7(prs)
    style_8(prs)
    style_9(prs)

    out = "presentation/banners/title_banners.pptx"
    prs.save(out)
    print(f"✓ Saved: {out}  (9 slides)")
    print("  1 — MACE doom block letters")
    print("  2 — ASCII IR spectrum art")
    print("  3 — system boot sequence")
    print("  4 — v1 spirit, minimal + centered")
    print("  5 — lean, thin strokes")
    print("  6 — lean, filled █")
    print("  7 — block, thin strokes")
    print("  8 — block, filled █")
    print("  9 — block + GAUSSIAN, centered")


if __name__ == "__main__":
    main()
