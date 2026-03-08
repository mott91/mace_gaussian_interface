#!/usr/bin/env python3
"""
MACE-Gaussian Interface — Research Presentation v2
March 2026 · ~15 minutes · 13 slides
Run: python generate_pptx_v2.py
Output: presentation_v2.pptx
"""

import json
import os

import pyfiglet
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Color scheme (GitHub Dark) ────────────────────────────────────────────────
BG     = RGBColor(13,  17,  23)   # #0D1117
TEXT   = RGBColor(201, 209, 217)  # #C9D1D9
ACCENT = RGBColor(88,  166, 255)  # #58A6FF  bright blue
GREEN  = RGBColor(63,  185, 80)   # #3FB950
YELLOW = RGBColor(210, 153, 34)   # #D29922
DIM    = RGBColor(139, 148, 158)  # #8B949E
RED    = RGBColor(248, 81,  73)   # #F85149
FONT   = "Consolas"

DATE   = "Manuel Ott · Hofer Lab · 2026-03-11"


# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def add_footer(slide, text=DATE):
    box = slide.shapes.add_textbox(Inches(1), Inches(6.85), Inches(8), Inches(0.4))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.name = FONT
    p.font.color.rgb = DIM
    p.alignment = PP_ALIGN.CENTER


def add_prompt(slide, command, y=0.35):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.6))
    tf = box.text_frame
    tf.text = command
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True


def content_slide(prs, command, lines):
    """Standard slide: command header + list of (text, color) lines."""
    slide = blank_slide(prs)
    add_prompt(slide, command)

    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.15), Inches(8.5), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True

    for i, (line, color) in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(4)

    add_footer(slide)
    return slide


def two_col_slide(prs, command, left_lines, right_lines, split=4.3):
    """Slide with two column panels."""
    slide = blank_slide(prs)
    add_prompt(slide, command)

    lbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(split - 0.3), Inches(5.4))
    ltf = lbox.text_frame
    ltf.word_wrap = True
    for i, (line, color) in enumerate(left_lines):
        p = ltf.paragraphs[i] if i == 0 else ltf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(3)

    rbox = slide.shapes.add_textbox(Inches(split), Inches(1.15), Inches(9.5 - split), Inches(5.4))
    rtf = rbox.text_frame
    rtf.word_wrap = True
    for i, (line, color) in enumerate(right_lines):
        p = rtf.paragraphs[i] if i == 0 else rtf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(3)

    add_footer(slide)
    return slide


# ── Data helpers ──────────────────────────────────────────────────────────────

def load_combo_metrics(molecule):
    """Load combo results from analysis_results_harmonic/{molecule}/data/metrics_summary.json.

    Data source: pre-aggregated metrics_summary.json (not comparison_results/ directory walk).
    Sorted by mae_freq ascending so the best-performing combo (lowest error) appears first.
    Per CONTEXT.md line 28: overrides original comparison_results/ approach from discuss-phase.
    """
    summary_path = os.path.join(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
        "analysis_results_harmonic",
        molecule,
        "data",
        "metrics_summary.json",
    )
    if not os.path.isfile(summary_path):
        return []
    with open(summary_path) as f:
        data = json.load(f)
    combos = []
    for entry in data.get("comparisons", []):
        combos.append({
            "name": entry.get("name", "unknown"),
            "mae_freq": entry.get("mae_freq", float("inf")),
            "r2_freq": entry.get("r2_freq", 0.0),
            "r2_intensity": entry.get("r2_intensity", 0.0),
            "speedup": entry.get("speedup"),
        })
    # Sort by mae_freq ascending (lowest MAE = best = shown first)
    # Per CONTEXT.md line 26: overrides original "R² descending" from discuss-phase.
    combos.sort(key=lambda c: c["mae_freq"], reverse=False)
    return combos


def parse_combo_name(name):
    """Split combo name into (energy_model, dipole_model)."""
    for suffix in ("_mace_ml", "_espaloma"):
        if name.endswith(suffix):
            return name[: -len(suffix)], suffix[1:]
    return name, "?"


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 0: style 9 — block font MACE + GAUSSIAN, centered."""
    slide = blank_slide(prs)

    # MACE in block font
    raw = pyfiglet.figlet_format("MACE", font="block")
    art_lines = [(ln, ACCENT) for ln in raw.splitlines() if ln.strip()]
    art_lines.append(("              GAUSSIAN", GREEN))

    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(9.2), Inches(3.0))
    tf = box.text_frame
    tf.word_wrap = False
    for i, (text, color) in enumerate(art_lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(0)
        p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(8.4), Inches(0.6))
    stf = sub.text_frame
    stf.text = "// ML-accelerated IR spectrum calculations"
    sp = stf.paragraphs[0]
    sp.font.size = Pt(16)
    sp.font.name = FONT
    sp.font.color.rgb = DIM
    sp.alignment = PP_ALIGN.CENTER

    add_footer(slide)


def slide_motivation(prs):
    """Slide 1: Why we built this."""
    content_slide(prs, "$ cat motivation.md", [
        ("# Problem", ACCENT),
        ("  \u2192 DFT frequency calculations: slow", TEXT),
        ("  \u2192 Minutes to hours per molecule", TEXT),
        ("  \u2192 Bottleneck for systematic benchmark campaigns", TEXT),
        ("", DIM),
        ("# Solution", ACCENT),
        ("  \u2192 Hybrid ML-QM approach", TEXT),
        ("  \u2192 ML dipoles and energies", TEXT),
        ("  \u2192 Possibly 10\u2013100\u00d7 faster than pure DFT", TEXT),
        ("", DIM),
        ("# Impact", ACCENT),
        ("  \u2192 Enable rapid IR spectral predictions", TEXT),
        ("  \u2192 Maintain DFT-level accuracy", TEXT),
        ("  \u2192 Systematic comparison: energy surface vs dipole model quality", TEXT),
    ])


def slide_ir_theory(prs):
    """Slide 2: IR spectroscopy — what Gaussian needs and why."""
    content_slide(prs, "$ cat ir_spectroscopy.md", [
        ("# How IR spectra are computed", ACCENT),
        ("  → Displace each atom in x, y, z  (small nudges around equilibrium)", TEXT),
        ("  → For each displacement, compute energy + forces", TEXT),
        ("  → Build Hessian (matrix of force constants) → diagonalise", TEXT),
        ("  → Eigenvalues = frequencies, eigenvectors = normal modes", TEXT),
        ("  → 3N−6 modes for N atoms  (e.g. water: 3 modes)", TEXT),
        ("", DIM),
        ("# Where do intensities come from?", ACCENT),
        ("  → Dipole μ⃗ = charge distribution of the molecule", TEXT),
        ("  → When atoms vibrate, μ⃗ changes → molecule absorbs IR light", TEXT),
        ("  → Intensity ∝ |∂μ/∂Q|²  (how much dipole changes per mode)", TEXT),
        ("  → So we also need μ⃗ at every displaced geometry", TEXT),
        ("", DIM),
        ("# The bottleneck", ACCENT),
        ("  → Gaussian displaces the molecule hundreds of times", TEXT),
        ("  → Each time: needs energy E, forces F, and dipole μ⃗", TEXT),
        ("  → With DFT: minutes to hours per molecule", TEXT),
        ("  → With ML:  seconds  ← this is what we replace", GREEN),
    ])


def slide_architecture(prs):
    slide = blank_slide(prs)
    add_prompt(slide, "$ python3 workflow.py")

    lines = [
        ("                              [molecule.xyz]", GREEN),
        ("                                    ▼", DIM),
        ("  ┌────────────────────────────────────────────────────────────────────────┐", DIM),
        ("  │  GEOMETRY OPTIMIZER   • load_geometry()  • mace_omol.optimize()        │", ACCENT),
        ("  └───────────────────────────────┬────────────────────────────────────────┘", DIM),
        ("                 ┌────────────────┴─────────────────┐", DIM),
        ("                 ▼                                  ▼", DIM),
        ("  ┌────────────────────────────────┐  ┌────────────────────────────────────┐", DIM),
        ("  │  DFT BASELINE                  │  │  ML FREQ CALC                      │", ACCENT),
        ("  │  ┌──────────────────┐          │  │  • load_mace_calculator()           │", DIM),
        ("  │  │ Built-in opt     │          │  │    (MACE-MP / MACE-OMOL)            │", DIM),
        ("  │  │ B3LYP/6-31G(d,p)│          │  │  • get_dipole_calculator()           │", DIM),
        ("  │  └──────────────────┘          │  │    (Espaloma / MACE_DIPOLE)          │", DIM),
        ("  │  • gaussian_freq()             │  │  • setup_zmq_server()               │", DIM),
        ("  │  • gaussian_dipoles()          │  │  • launch_gaussian()                │", DIM),
        ("  │  • parse_gaussian()            │  │  • zmq_dipole_loop()                │", DIM),
        ("  └──────────────┬─────────────────┘  └──────────────────┬─────────────────┘", DIM),
        ("                 └────────────────┬─────────────────────┘", DIM),
        ("                                  ▼", DIM),
        ("  ┌────────────────────────────────────────────────────────────────────────┐", DIM),
        ("  │  ANALYSIS   load_results → eigenvector_matching → kde → metrics → plot │", ACCENT),
        ("  └───────────────────────────────┬────────────────────────────────────────┘", DIM),
        ("                                  ▼", DIM),
        ("                        generate_html_report()", GREEN),
    ]

    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.15), Inches(8.5), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = False
    for i, (text, color) in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(0)

    add_footer(slide)


def slide_dipole_methods(prs):
    two_col_slide(prs, "$ ls -l calculators/",
        left_lines=[
            ("# MACE4IR (direct dipole)", ACCENT),
            ("", DIM),
            ("  Input: atomic positions + species", TEXT),
            ("       ↓  E(3)-equivariant GNN", TEXT),
            ("  Output: μ = (μx, μy, μz)", TEXT),
            ("", DIM),
            ("  → Equivariant: respects rotations", TEXT),
            ("  → Trained specifically on dipoles", TEXT),
            ("  → 78-element coverage", TEXT),
            ("  → Custom MACE4IR model (78 el.)", TEXT),
            ("  → Fast — single forward pass", TEXT),
            ("", DIM),
            ("# Why E(3)-equivariance matters", ACCENT),
            ("  Dipole is a vector — it must rotate", TEXT),
            ("  with the molecule. Standard NNs", TEXT),
            ("  don't guarantee this.", TEXT),
        ],
        right_lines=[
            ("# Espaloma (charge-based dipole)", ACCENT),
            ("", DIM),
            ("  Input: molecular graph", TEXT),
            ("       ↓  graph neural network", TEXT),
            ("  Output: partial charges q₁…qₙ", TEXT),
            ("       ↓  μ = Σ qᵢ · rᵢ", TEXT),
            ("  Output: dipole vector", TEXT),
            ("", DIM),
            ("  → Physics-based construction", TEXT),
            ("  → Interpretable (charges visible)", TEXT),
            ("  → Available as pip package", TEXT),
            ("  → No custom model needed", TEXT),
            ("", DIM),
            ("# Energy calculators (PES)", ACCENT),
            ("", DIM),
            ("  MACE-OMOL-0   molecules, default", GREEN),
            ("  MACE-MP-0     universal (130 el.)", GREEN),
            ("  MACE-OFF      organic molecules", GREEN),
            ("  MACE-ANI-CC   organic (CC training data)", GREEN),
            ("", DIM),
            ("  → 4 energy × 2 dipole = 8 combos", TEXT),
        ],
        split=5.1,
    )


def slide_zmq(prs):
    """Slide 5: ZMQ bridge — ASCII diagram + External keyword."""
    slide = blank_slide(prs)
    add_prompt(slide, "$ cat zmq_bridge.md")

    # Brief refresher + External keyword
    top_lines = [
        ("# Gaussian needs E, F, μ⃗ at every displaced geometry — how do we inject ML?", ACCENT),
        ("  → Gaussian 'External' keyword: calls an external script per displacement", TEXT),
        ("  → We intercept this with a ZMQ bridge — no Gaussian modifications needed", TEXT),
    ]

    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9.0), Inches(1.2))
    ttf = tbox.text_frame
    ttf.word_wrap = True
    for i, (text, color) in enumerate(top_lines):
        p = ttf.paragraphs[i] if i == 0 else ttf.add_paragraph()
        p.text = text
        p.font.size = Pt(13)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(2)

    # ASCII flow diagram
    diagram = [
        ("  ┌─────────────────────┐                    ┌─────────────────────────────┐", DIM),
        ("  │                     │   displacement #n  │                             │", DIM),
        ("  │    GAUSSIAN 16      │ ──────────────────→│     gm_helper.py            │", ACCENT),
        ("  │                     │                    │     (External script)        │", DIM),
        ("  │    freq=anharmonic  │                    └──────────────┬──────────────┘", DIM),
        ("  │    external='...'   │                                   │ ZMQ IPC", DIM),
        ("  │                     │                                   ▼", DIM),
        ("  │                     │                    ┌─────────────────────────────┐", DIM),
        ("  │                     │                    │     zmq_server.py           │", ACCENT),
        ("  │                     │                    │                             │", DIM),
        ("  │                     │     fort.7         │  ┌─ MACE energy → E, F     │", GREEN),
        ("  │                     │ ←──────────────────│  └─ MACE dipole → μ⃗        │", GREEN),
        ("  │                     │   E, F, μ⃗          │                             │", DIM),
        ("  └─────────────────────┘                    └─────────────────────────────┘", DIM),
        ("", DIM),
        ("  ↻ repeats for every displacement (hundreds per molecule)", TEXT),
    ]

    dbox = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(8.5), Inches(4.0))
    dtf = dbox.text_frame
    dtf.word_wrap = False
    for i, (text, color) in enumerate(diagram):
        p = dtf.paragraphs[i] if i == 0 else dtf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.font.name = FONT
        p.font.color.rgb = color
        p.space_after = Pt(0)

    add_footer(slide)


def slide_mode_matching(prs):
    content_slide(prs, "$ cat mode_matching.md", [
        ("# Problem: ML and DFT modes can be in different order", ACCENT),
        ("", DIM),
        ("  DFT:  [1595, 1610, 2900, 3050, 3100] cm⁻¹", TEXT),
        ("  ML:   [1608, 1590, 2905, 3048, 3102] cm⁻¹   ← modes 1 & 2 swapped!", YELLOW),
        ("", DIM),
        ("  → 1595 and 1610 are close but physically different modes", TEXT),
        ("", DIM),
        ("# Solution: Eigenvector dot product", ACCENT),
        ("  Each mode = eigenvector (3N-dim displacement pattern)", TEXT),
        ("  Same physical motion → parallel eigenvectors", TEXT),
        ("  similarity(i, j) = |v_ML[i] · v_DFT[j]|     (0 = orthogonal, 1 = same)", TEXT),
        ("", DIM),
        ("  → Match by motion, not by frequency", TEXT),
        ("", DIM),
        ("# Example", ACCENT),
        ("  ML 1608 cm⁻¹ → DFT 1595 cm⁻¹  [overlap: 0.12]   ✗ closest freq, wrong mode", RED),
        ("  ML 1608 cm⁻¹ → DFT 1610 cm⁻¹  [overlap: 0.97]   ✓ correct — same motion", GREEN),
    ])


def slide_results_overview(prs):
    """Slide 7: HTML report listing with ASCII molecule art."""
    content_slide(prs, "$ ls analysis_results_harmonic/", [
        ("# Spectrum analysis reports", ACCENT),
        ("", DIM),
        ("       O                       O   OH", GREEN),
        ("      / \\                      ‖   |", GREEN),
        ("     H   H   H₂O          CH₃-C-O-C₆H₄-COOH", GREEN),
        ("                               aspirin C₉H₈O₄", DIM),
        ("", DIM),
        ("  water/report.html       3 atoms  · 3 modes  · 8 combos", GREEN),
        ("  aspirin/report.html    19 atoms  · 51 modes · 4 combos", GREEN),
        ("", DIM),
        ("  Each report: regression plots · KDE spectra · mode matching", TEXT),
        ("  Open: $ open analysis_results_harmonic/water/report.html", DIM),
        ("         $ open analysis_results_harmonic/aspirin/report.html", DIM),
    ])


def slide_results_table(prs):
    """Slides 8–9: Harmonic benchmark — water (8 combos) and aspirin (4 combos).

    Rows sorted by MAE(freq) ascending (lowest error = best = shown first).
    Uniform TEXT color — no per-row color logic.
    Data loaded from analysis_results_harmonic/{molecule}/data/metrics_summary.json.
    """
    water = load_combo_metrics("water")
    aspirin = load_combo_metrics("aspirin")

    def table_lines(combos, mol_label):
        lines = [
            (f"# Harmonic benchmark — {mol_label}", ACCENT),
            ("  Ref: B3LYP/6-31G(d,p)   |   sorted by MAE(freq) ascending", DIM),
            ("", DIM),
            ("  Energy model    Dipole      MAE(freq)   R²(freq)    R²(intens)", DIM),
            ("  " + "─" * 60, DIM),
        ]
        for c in combos:
            energy, dipole = parse_combo_name(c["name"])
            row = (
                f"  {energy:<14}  {dipole:<10}"
                f"  {c['mae_freq']:>6.1f} cm⁻¹"
                f"   {c['r2_freq']:.7f}"
                f"   {c['r2_intensity']:.2f}"
            )
            lines.append((row, TEXT))
        return lines

    water_lines = table_lines(water, "water (H₂O, 3 modes)")
    aspirin_lines = table_lines(aspirin, "aspirin (C₉H₈O₄, 51 modes)")

    # Water slide
    content_slide(prs, "$ cat analysis_results_harmonic/water/data/metrics_summary.json | sort-by-mae", water_lines + [
        ("", DIM),
        ("  → Frequencies: all combos excellent (R² > 0.9999)", TEXT),
        ("  → Intensities: mace_ml consistently beats espaloma", TEXT),
    ])

    # Aspirin slide
    content_slide(prs, "$ cat analysis_results_harmonic/aspirin/data/metrics_summary.json | sort-by-mae", aspirin_lines + [
        ("", DIM),
        ("  → mace_omol: best accuracy (MAE < 9 cm⁻¹, R² 0.9998)", TEXT),
        ("  → mace_mp: most speedup (~29×) but worse accuracy", TEXT),
    ])


def slide_scaling(prs):
    """Slide 10: Molecule size vs speedup — ASCII bar chart."""
    MAX_BAR = 20
    molecules = [
        ("water",   3,  1,  "~1×"),
        ("glycine", 10, 6,  "~7–10×"),
        ("aspirin", 21, 20, "~18–29×"),
    ]
    bar_lines = []
    for name, atoms, bar_val, label in molecules:
        bar = "=" * bar_val + " " * (MAX_BAR - bar_val)
        bar_lines.append((
            f"  {name:<10}  ({atoms:>2} atoms)  |{bar}|  {label}",
            TEXT,
        ))

    content_slide(prs, "$ python benchmark.py --scaling", [
        ("# ML speedup scales with molecule size", ACCENT),
        ("  Dipole step: ML vs DFT (B3LYP/6-31G(d,p))", DIM),
        ("", DIM),
        ("  molecule      atoms         speedup", DIM),
        ("  " + "─" * 50, DIM),
        *bar_lines,
        ("", DIM),
        ("  → Water (3 atoms): ML overhead visible — no speedup at this scale", TEXT),
        ("  → Glycine (10 atoms): ~7–10× — target use case", TEXT),
        ("  → Aspirin (21 atoms): ~18–29× — strong benefit", TEXT),
        ("", DIM),
        ("  Speedup source: DFT dipoles O(N³); ML dipoles O(N)", DIM),
        ("  (estimated ranges — not benchmarked to the second)", DIM),
    ])


def slide_status(prs):
    """Slide 9: Research journey — what we built, ran, found, and what's next."""
    two_col_slide(prs, "$ git log --oneline --graph",
        left_lines=[
            ("# What we built", ACCENT),
            ("  ZMQ bridge: ML dipoles \u2192 Gaussian VPT2", TEXT),
            ("  CLI: mace-gaussian run molecule.xyz", DIM),
            ("  4 energy \u00d7 2 dipole models = 8 combinations", TEXT),
            ("", DIM),
            ("# What we ran", ACCENT),
            ("  Harmonic benchmark across 8 molecules", TEXT),
            ("  Water, aspirin, glycine, methane, ethane, NH\u2083, CO", DIM),
            ("", DIM),
            ("# What we found", ACCENT),
            ("  Frequencies: R\u00b2 > 0.999 across all combos \u2713", GREEN),
            ("  Intensities: mace_ml >> espaloma (0.72 vs 0.52)", GREEN),
            ("  Speedup scales with molecule size (aspirin: ~18\u00d7)", GREEN),
        ],
        right_lines=[
            ("# Still open", ACCENT),
            ("", DIM),
            ("  Anharmonic validation:", TEXT),
            ("  \u2192 VPT2 pipeline runs; accuracy TBD", DIM),
            ("  \u2192 Overtones + combination bands", DIM),
            ("", DIM),
            ("  25-molecule benchmark campaign:", TEXT),
            ("  \u2192 Systematic across functional groups", DIM),
            ("  \u2192 Batch runner in progress", DIM),
            ("", DIM),
            ("# Thesis question", ACCENT),
            ("  Energy surface quality", TEXT),
            ("  vs dipole model quality \u2014", TEXT),
            ("  which dominates IR accuracy?", TEXT),
            ("", DIM),
            ("  $ mace-gaussian run molecule.xyz", GREEN),
        ],
        split=5.1,
    )


def slide_questions(prs):
    slide = blank_slide(prs)

    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.0))
    tf = box.text_frame

    p = tf.paragraphs[0]
    p.text = "$ ./questions.sh"
    p.font.size = Pt(36)
    p.font.name = FONT
    p.font.color.rgb = ACCENT
    p.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = ""

    p3 = tf.add_paragraph()
    p3.text = "// Thank you"
    p3.font.size = Pt(18)
    p3.font.name = FONT
    p3.font.color.rgb = DIM

    add_footer(slide)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)            # 0  — neofetch header
    slide_motivation(prs)       # 1  — problem / solution / impact
    slide_ir_theory(prs)        # 2  — what is IR, harmonic limit
    slide_architecture(prs)     # 3  — system architecture diagram
    slide_dipole_methods(prs)   # 4  — MACE4IR + Espaloma + energy models
    slide_zmq(prs)              # 5  — ZMQ bridge — pedagogical flow
    slide_mode_matching(prs)    # 6  — eigenvector dot product
    slide_results_overview(prs) # 7  — HTML report listing + ASCII art
    slide_results_table(prs)    # 8–9 — harmonic table (water + aspirin, 2 slides)
    slide_scaling(prs)          # 10 — molecule size vs speedup bar chart
    slide_status(prs)           # 11 — research journey
    slide_questions(prs)        # 12 — questions

    out = "presentation/presentation_v2.pptx"
    prs.save(out)

    print(f"✓ Saved: {out}")
    print(f"  13 slides  (~15 min)")
    print(f"  Slides: title · motivation · IR theory ·")
    print(f"          architecture · calculators · ZMQ bridge ·")
    print(f"          mode matching · results overview ·")
    print(f"          results table (water) · results table (aspirin) ·")
    print(f"          scaling · status · questions")


if __name__ == "__main__":
    main()
